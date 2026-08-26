"""
Rapports Word (.docx), PowerPoint (.pptx) et PDF du menu Générateur de rapport.

Deux principes suivis strictement, en réponse directe aux défauts observés :

  1. Rien n'est laissé de côté. Le tableau de synthèse et la section
     détaillée couvrent TOUTES les analyses de la session — y compris
     celles qui n'ont pas de graphique associé (ex : test de lien entre
     deux facteurs) — jamais seulement celles qui ont une image.
     L'interprétation de chaque analyse réutilise directement
     `viz.comments.generate_comment_rules`, qui couvre déjà tous les types
     de résultats produits par l'application (categorical, numeric,
     grouped_numeric, timeseries, scatter, correlation, rate, group_test,
     association_test, period_comparison) — pas de type oublié.

  2. Le PowerPoint n'essaie jamais de tout charger sur une seule diapositive.
     Le texte du rapport et le tableau de synthèse sont automatiquement
     répartis sur autant de diapositives que nécessaire (jamais de
     troncature silencieuse), et chaque graphique a sa propre diapositive.

Aucune donnée brute n'est utilisée ici : uniquement du texte déjà généré et
des images de graphiques déjà tracés côté application.
"""
from __future__ import annotations

import io
import re
from datetime import datetime
from pathlib import Path

ORANGE = "F28C28"
# --- Règles de présentation d'un rapport statistique professionnel ---------
# Appliquées aux 4 générateurs "rapport" actifs (PDF/Word × questionnaire/
# fichier brut) — pas aux exports PowerPoint (diapositives, conventions de
# mise en page différentes) ni aux fonctions historiques non utilisées par
# l'interface actuelle. Couleurs volontairement INCHANGÉES (déjà validées).
REPORT_MARGIN_TOP_CM = 2.5
REPORT_MARGIN_BOTTOM_CM = 2.5
REPORT_MARGIN_LEFT_CM = 3.0
REPORT_MARGIN_RIGHT_CM = 2.5
REPORT_FONT = "Times-Roman"       # Times New Roman (police standard PDF, pas d'empaquetage nécessaire)
REPORT_FONT_BOLD = "Times-Bold"
REPORT_FONT_ITALIC = "Times-Italic"
REPORT_BODY_SIZE = 12
REPORT_LINE_SPACING = 1.5         # interligne 1,5 — leading = taille × 1,5
REPORT_TITLE_SIZE = 15            # 14-16pt demandé pour les titres, gras
# Palette pour les graphiques à plusieurs séries (croisements qualitatif ×
# qualitatif, comparaisons) — mêmes couleurs que la palette JS du Tableau
# de bord générique (assets/dashboard_generic_auto.html, const PALETTE),
# pour qu'une série garde la même couleur à l'écran et dans le rapport.
PALETTE_HEX = ["#3366cc", "#dc3912", "#ff9900", "#109618", "#990099", "#0099c6", "#dd4477", "#66aa00"]
YELLOW = "FFD54F"
WHITE = "FFFFFF"
BLACK = "000000"


def _fmt_pct(v) -> str:
    return f"{v:.2f} %" if v is not None else "N/A"


def _fmt_score(v) -> str:
    return f"{v:.2f}/5" if v is not None else "N/A"


def _fmt_min(v) -> str:
    return f"{v:.1f} min" if v is not None else "N/A"


def _fmt_int(v) -> str:
    return str(v) if v is not None else "N/A"


def _fmt_nps(v) -> str:
    # NPS peut être négatif — signe "+" explicite pour les valeurs
    # positives (convention standard d'affichage du NPS), pas de "%" (ce
    # n'est pas un pourcentage, l'échelle va de -100 à +100).
    if v is None:
        return "N/A"
    return f"+{v}" if v > 0 else str(v)


# Pastilles KPI supplémentaires (au-delà des 5 d'origine du gabarit PPTX),
# proposables au choix dans le Générateur de rapport. Même structure que les
# pastilles d'origine : (clé, libellé, formateur, couleur de fond, couleur
# de texte) — partagé entre PDF, Word et PowerPoint pour ne pas dupliquer 3
# fois le même style.
EXTRA_KPI_PILLS = [
    ("csat", "CSAT", _fmt_pct, "FFE082", "1a1a1a"),
    ("ces_estime", "CES estimé", _fmt_pct, "E0E0E0", "1a1a1a"),
    ("nps_estime", "NPS estimé", _fmt_nps, "B2DFDB", "1a1a1a"),
    ("duree_moyenne_min", "Durée Moyenne", _fmt_min, "D1C4E9", "1a1a1a"),
    ("n_telephones", "Contacts Tél.", _fmt_int, "FFCCBC", "1a1a1a"),
]

# Nombre maximal de lignes / caractères de texte par diapositive PowerPoint,
# pour ne jamais saturer une page — au-delà, on ouvre une diapositive
# supplémentaire plutôt que de tasser le contenu.
PPTX_MAX_LINES_PER_SLIDE = 8
PPTX_MAX_CHARS_PER_SLIDE = 650
PPTX_MAX_TABLE_ROWS_PER_SLIDE = 10


def _strip_emoji(text: str) -> str:
    """Retire les émojis d'un libellé avant de l'envoyer à matplotlib : la
    police standard (DejaVu Sans) ne les a pas, ils s'affichaient comme des
    carrés vides dans le PDF/Word."""
    return re.sub(r"[\U0001F300-\U0001FAFF\U00002600-\U000027BF\U0001F1E6-\U0001F1FF]+\s*", "", str(text)).strip()


def _donut_png_matplotlib(labels: list[str], values: list[float], colors_hex: list[str],
                           width_px: int = 700, height_px: int = 450) -> bytes | None:
    """Anneau de satisfaction en matplotlib — remplace le rendu Plotly/kaleido
    pour le PDF/Word du modèle officiel. Pur CPU (pas de mini-navigateur en
    tâche de fond) : nettement plus léger que kaleido sur un hébergement aux
    ressources limitées, où kaleido peut devenir très lent, voire figer,
    surtout répété sur plusieurs agences d'affilée."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    try:
        dpi = 100
        fig, ax = plt.subplots(figsize=(width_px / dpi, height_px / dpi), dpi=dpi)
        colors = [f"#{c.lstrip('#')}" for c in colors_hex] if colors_hex else None
        clean_labels = [_strip_emoji(l) for l in labels]
        wedges, _ = ax.pie(values, colors=colors, startangle=90, wedgeprops={"width": 0.42, "edgecolor": "white"})
        ax.legend(wedges, [f"{l} ({v})" for l, v in zip(clean_labels, values)], loc="center left",
                  bbox_to_anchor=(1.0, 0.5), fontsize=9, frameon=False)
        ax.set_aspect("equal")
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return None


def _bar_png_matplotlib(labels: list[str], values: list[float],
                         width_px: int = 700, height_px: int = 450) -> bytes | None:
    """Barres horizontales (points appréciés) en matplotlib — même raison
    d'être que `_donut_png_matplotlib` ci-dessus."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    try:
        dpi = 100
        fig, ax = plt.subplots(figsize=(width_px / dpi, height_px / dpi), dpi=dpi)
        y_pos = range(len(labels))
        ax.barh(list(y_pos), values, color=f"#{ORANGE}")
        ax.set_yticks(list(y_pos))
        ax.set_yticklabels(labels, fontsize=9)
        ax.invert_yaxis()
        ax.set_xlabel("%", fontsize=9)
        for i, v in enumerate(values):
            ax.text(v + 1, i, f"{v}%", va="center", fontsize=8.5)
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return None


def _agency_chart_png(kind: str, page: dict, width_px: int = 700, height_px: int = 450) -> bytes | None:
    """Ancien chemin de rendu (donut/bar fixes) — plus utilisé depuis que
    `build_agency_full_charts` fournit TOUS les graphiques d'une agence sans
    restriction (voir `page["all_charts"]`), rendus via `figure_to_png_bytes`.
    Conservée uniquement si un appelant externe s'y référait encore."""
    cache_key = f"_{kind}_png"
    if page.get(cache_key) is not None:
        return page[cache_key]
    fig = page.get(kind)
    result = figure_to_png_bytes(fig, width=width_px, height=height_px, scale=1.5) if fig is not None else None
    if result is not None:
        page[cache_key] = result
    return result


def _generic_plotly_to_png_matplotlib(fig, width_px: int = 1200, height_px: int = 650, show_title: bool = True) -> bytes | None:
    """Convertit une figure Plotly « simple » (Bar, Scatter/ligne, ou Pie —
    les seuls types produits par `viz.report_charts.build_report_charts`)
    en PNG via matplotlib, sans passer par kaleido. Retourne None si la
    figure ne correspond à aucun de ces cas simples, pour que l'appelant
    retombe alors sur kaleido plutôt que de perdre le graphique — cette
    fonction ne couvre que les cas réels de la galerie, pas un convertisseur
    Plotly générique complet."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None

    try:
        traces = list(fig.data)
        if not traces or not all(t.type == traces[0].type for t in traces):
            return None
        ttype = traces[0].type
        # Un seul type de figure gère plusieurs traces : les barres groupées
        # multi-séries (ex : "Taux de satisfaction / insatisfaction par
        # agence") ET les boîtes à moustaches (une trace par catégorie).
        # Pour tout le reste, une seule trace est attendue.
        if len(traces) > 1 and ttype not in ("bar", "box"):
            return None
        trace = traces[0]
        title = ""
        if fig.layout.title and fig.layout.title.text:
            title = re.sub(r"<[^>]+>", "", fig.layout.title.text)

        # Tailles de police alignées sur `viz.charts.apply_readable_style`
        # (celles utilisées à l'écran), puis encore agrandies pour rester
        # lisibles une fois imprimées à taille réelle sur une page A4 —
        # les valeurs précédentes (14-20pt) paraissaient minuscules une fois
        # le graphique inséré dans le rapport final.
        FS_TITLE = 26
        FS_TICK = 19
        FS_LABEL = 18
        FS_LEGEND = 17

        dpi = 100
        mpl_fig, ax = plt.subplots(figsize=(width_px / dpi, height_px / dpi), dpi=dpi)
        # Un seul cadre au total (demande explicite : "il ne doit pas y
        # avoir plusieurs cadres... juste les axes et le cadre") — sans
        # ceci, la boîte par défaut de matplotlib (4 bordures autour de la
        # seule zone de tracé) se superposait au cadre ajouté plus bas
        # autour de TOUTE la figure (qui, lui, englobe aussi les étiquettes
        # en biais de l'axe des x) : deux rectangles visibles au lieu d'un
        # seul. Ne restent que les deux axes naturels (bas, gauche) ; le
        # rectangle est refermé par le cadre extérieur, pas par une
        # seconde boîte intérieure.
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

        if ttype == "bar" and len(traces) > 1:
            # Barres groupées : une catégorie par valeur de x, une couleur
            # par trace (série), décalées côte à côte.
            import numpy as _np
            xs = [str(v) for v in (traces[0].x or [])]
            n_series = len(traces)
            bar_width = 0.8 / n_series
            positions = _np.arange(len(xs))
            for i, tr in enumerate(traces):
                ys = list(tr.y) if tr.y is not None else []
                color = tr.marker.color if tr.marker and tr.marker.color is not None else f"#{ORANGE}"
                if isinstance(color, (list, tuple)):
                    color = str(color[0]) if color else f"#{ORANGE}"
                offset = (i - (n_series - 1) / 2) * bar_width
                bars = ax.bar(positions + offset, ys, width=bar_width, label=tr.name or f"Série {i+1}", color=color)
                # Utilise le texte personnalisé du trace (pourcentages,
                # défini côté report_charts.py) au lieu de la valeur brute
                # par défaut — bug réel trouvé : ax.bar_label() sans
                # "labels=" explicite ignore totalement trace.text et
                # réaffiche la valeur numérique brute de chaque barre.
                bar_texts = list(tr.text) if tr.text else None
                ax.bar_label(bars, labels=bar_texts, fontsize=FS_LABEL - 2, padding=2)
            ax.set_xticks(positions)
            ax.set_xticklabels([_strip_emoji(x) for x in xs], rotation=20, ha="right", fontsize=FS_TICK)
            ax.tick_params(axis="y", labelsize=FS_TICK)
            ax.legend(fontsize=FS_LEGEND, frameon=False, loc="upper center",
                      bbox_to_anchor=(0.5, -0.18), ncol=n_series)
        elif ttype == "bar":
            xs = list(trace.x) if trace.x is not None else []
            ys = list(trace.y) if trace.y is not None else []
            orientation = getattr(trace, "orientation", None)
            colors = trace.marker.color if trace.marker and trace.marker.color is not None else f"#{ORANGE}"
            if isinstance(colors, (list, tuple)):
                colors = [str(c) for c in colors]
            if orientation == "h":
                # Convention Plotly pour les barres horizontales : x = valeurs,
                # y = catégories (inverse du cas vertical) — d'où l'inversion
                # ici, sinon les libellés de catégorie se retrouvaient comme
                # position numérique et les valeurs comme libellés (bug corrigé).
                categories_clean = [_strip_emoji(v) for v in ys]
                # Positions de graduation FIXÉES explicitement (0, 1, 2...) au
                # lieu de laisser ax.barh() gérer l'axe des catégories tout
                # seul — bug réel trouvé et confirmé : avec UNE SEULE
                # catégorie, matplotlib générait plusieurs graduations
                # automatiques sur l'axe (numérique par défaut) et leur
                # appliquait TOUTES le même texte de catégorie, donnant
                # l'impression que le même libellé se répétait 6-8 fois.
                y_positions = list(range(len(categories_clean)))
                bars = ax.barh(y_positions, xs, color=colors, height=0.5)
                ax.set_yticks(y_positions)
                ax.set_yticklabels(categories_clean)
                ax.invert_yaxis()
                # Axe X (celui des VALEURS ici, barres horizontales) en
                # nombres entiers uniquement — sans ça, une seule catégorie
                # avec une petite valeur (ex: 1 citation) faisait générer à
                # matplotlib des graduations à virgule absurdes (0.2, 0.4...).
                from matplotlib.ticker import MaxNLocator as _MaxNLocatorH
                ax.xaxis.set_major_locator(_MaxNLocatorH(integer=True))
                # Espace vertical minimum réservé (comme s'il y avait au
                # moins 4 catégories) — demande explicite : "si on a une
                # seule modalité, je ne veux pas voir des grosses barres".
                # Sans ceci, avec une seule catégorie, matplotlib étire
                # l'axe Y pour occuper tout l'espace disponible et la barre
                # (même à hauteur fixe de 0.5) paraît énorme faute d'espace
                # vide autour — avec cette réserve, une seule barre garde la
                # même épaisseur visuelle qu'avec plusieurs.
                min_categories = 4
                if len(categories_clean) < min_categories:
                    pad = (min_categories - len(categories_clean)) / 2
                    ax.set_ylim(len(categories_clean) - 1 + pad + 0.5, -pad - 0.5)
                ax.bar_label(bars, labels=(list(trace.text) if trace.text else None), fontsize=FS_LABEL, padding=4)
                ax.tick_params(axis="y", labelsize=FS_TICK)
                ax.tick_params(axis="x", labelsize=FS_TICK)
                # Laisse la place aux étiquettes de catégorie longues
                # ("Le temps de traitement de ma demande"...), sans quoi
                # matplotlib les tronque ou les fait sortir du cadre.
                max_label_len = max((len(c) for c in categories_clean), default=0)
                mpl_fig.subplots_adjust(left=min(0.05 + 0.012 * max_label_len, 0.55))
            else:
                xs_clean = [_strip_emoji(x) for x in xs]
                bars = ax.bar(xs_clean, ys, color=colors)
                ax.bar_label(bars, labels=(list(trace.text) if trace.text else None), fontsize=FS_LABEL, padding=4)
                plt.setp(ax.get_xticklabels(), rotation=20, ha="right", fontsize=FS_TICK)
                ax.tick_params(axis="y", labelsize=FS_TICK)
        elif ttype == "scatter":
            xs = list(trace.x) if trace.x is not None else []
            ys = list(trace.y) if trace.y is not None else []
            # Un VRAI nuage de points (mode "markers" seul, croisement de 2
            # variables quantitatives — corrélation, pas une séquence) et
            # une COURBE (mode "lines+markers", évolution/tendance) sont
            # tous deux un trace Plotly "scatter" — seul `trace.mode` les
            # distingue. Avant, les deux étaient traités pareil : un nuage
            # de points se retrouvait avec des points reliés par un trait
            # (suggérant à tort un ordre/une séquence) ET une étiquette de
            # valeur sur CHAQUE point, illisible dès qu'il y a plus de
            # quelques dizaines de points.
            mode = getattr(trace, "mode", None) or "lines+markers"
            is_pure_scatter = "lines" not in mode
            if is_pure_scatter:
                ax.scatter(xs, ys, color=f"#{ORANGE}", s=70, alpha=.75, edgecolors="white", linewidths=.5)
            else:
                ax.plot(xs, ys, marker="o", color=f"#{ORANGE}", linewidth=2.5, markersize=7)
                for x, y in zip(xs, ys):
                    ax.annotate(str(y), (x, y), textcoords="offset points", xytext=(0, 8),
                                fontsize=FS_LABEL, ha="center")
            plt.setp(ax.get_xticklabels(), rotation=20, ha="right", fontsize=FS_TICK)
            ax.tick_params(axis="y", labelsize=FS_TICK)
        elif ttype == "pie":
            labels = list(trace.labels) if trace.labels is not None else []
            values = list(trace.values) if trace.values is not None else []
            colors_raw = trace.marker.colors if trace.marker and trace.marker.colors is not None else None
            hole = getattr(trace, "hole", 0) or 0
            total = sum(values) or 1
            # Étiquette directement sur chaque part (valeur + %) — sans ça,
            # seule la légende affichait un chiffre, ce qui rendait le
            # graphique illisible isolément (ex: sur une diapositive, ou
            # imprimé sans la légende à côté).
            def _slice_label(v):
                return f"{v}\n({100*v/total:.1f}%)"
            wedges, _texts, _autotexts = ax.pie(
                values, colors=colors_raw, startangle=90,
                wedgeprops={"width": 1 - hole, "edgecolor": "white"} if hole else {"edgecolor": "white"},
                autopct=lambda pct: _slice_label(round(pct * total / 100)),
                pctdistance=0.75 if hole else 0.6,
                textprops={"fontsize": FS_LABEL - 3, "fontweight": "bold", "color": "white" if not hole else "#1a1a1a"},
            )
            ax.legend(wedges, [f"{_strip_emoji(l)} ({v})" for l, v in zip(labels, values)],
                      loc="center left", bbox_to_anchor=(1.0, 0.5), fontsize=FS_LEGEND, frameon=False)
            ax.set_aspect("equal")
        elif ttype == "box":
            # Boîte à moustaches — deux cas possibles :
            #  1) une trace par catégorie, chacune avec les valeurs BRUTES
            #     (trace.y rempli) -> matplotlib calcule lui-même les
            #     quartiles, comme Plotly le fait à l'écran (ax.boxplot).
            #  2) une seule trace avec des quartiles DÉJÀ CALCULÉS
            #     (trace.q1/median/q3 remplis, trace.y vide) -> matplotlib
            #     doit recevoir ces statistiques toutes faites (ax.bxp),
            #     jamais recalculées à partir de rien.
            if traces[0].q1 is not None:
                tr = traces[0]
                xs = [str(v) for v in (tr.x or [])]
                stats = []
                for i, lab in enumerate(xs):
                    stats.append({
                        "label": lab, "med": tr.median[i], "q1": tr.q1[i], "q3": tr.q3[i],
                        "whislo": tr.lowerfence[i] if tr.lowerfence is not None else tr.q1[i],
                        "whishi": tr.upperfence[i] if tr.upperfence is not None else tr.q3[i],
                        "fliers": [],
                    })
                bp = ax.bxp(stats, patch_artist=True, showfliers=False)
                for box in bp["boxes"]:
                    box.set_facecolor(f"#{ORANGE}"); box.set_alpha(.55); box.set_edgecolor(f"#{ORANGE}")
                for med in bp["medians"]:
                    med.set_color("#1a1a1a"); med.set_linewidth(2)
                ax.tick_params(axis="x", labelsize=FS_TICK)
                ax.tick_params(axis="y", labelsize=FS_TICK)
                plt.setp(ax.get_xticklabels(), rotation=20, ha="right")
            else:
                data_per_box = [list(tr.y) if tr.y is not None else [] for tr in traces]
                box_labels = [_strip_emoji(tr.name or f"Boîte {i+1}") for i, tr in enumerate(traces)]
                bp = ax.boxplot(data_per_box, labels=box_labels, patch_artist=True, showfliers=True)
                for box in bp["boxes"]:
                    box.set_facecolor(f"#{ORANGE}"); box.set_alpha(.55); box.set_edgecolor(f"#{ORANGE}")
                for med in bp["medians"]:
                    med.set_color("#1a1a1a"); med.set_linewidth(2)
                ax.tick_params(axis="x", labelsize=FS_TICK)
                ax.tick_params(axis="y", labelsize=FS_TICK)
                plt.setp(ax.get_xticklabels(), rotation=20, ha="right")
            y_title = fig.layout.yaxis.title.text if fig.layout.yaxis and fig.layout.yaxis.title else None
            if y_title:
                ax.set_ylabel(_strip_emoji(y_title), fontsize=FS_LABEL, fontweight="bold")
        else:
            plt.close(mpl_fig)
            return None

        if title and show_title:
            ax.set_title(_strip_emoji(title), fontsize=FS_TITLE, fontweight="bold", pad=14)
        # Axes en nombres entiers uniquement (demande explicite : "pas de
        # valeurs décimales dans les axes... que des nombres entiers") —
        # sans ceci, matplotlib choisit parfois des graduations à virgule
        # (0.5, 1.5...) selon l'amplitude des données. Uniquement l'axe Y
        # ici : l'axe X peut porter une vraie variable continue (nuage de
        # points, nuage de valeurs) où forcer des entiers casserait la
        # précision réelle des données — voir plus haut pour le cas des
        # barres HORIZONTALES, où l'axe des valeurs (X) est traité à part.
        from matplotlib.ticker import MaxNLocator
        ax.yaxis.set_major_locator(MaxNLocator(integer=True))
        # Un peu d'espace entre le contenu (barres, étiquettes) et le cadre
        # — demande explicite : "trop collé, mets de l'espace, on doit voir
        # clairement". Marge interne entre les données et les bords de la
        # zone de tracé (mpl_fig.tight_layout(), appelé juste après,
        # calcule déjà automatiquement les marges externes — ajouter un
        # subplots_adjust ici serait aussitôt écrasé).
        ax.margins(y=0.08)
        # Cadre visible autour de TOUTE la figure (pas juste la zone de
        # tracé) — demande explicite : "les cadres du graphique, laisse les
        # axes à l'intérieur... tout ça doit être à l'intérieur du cadre".
        # fig.patch couvre l'intégralité de l'image (coordonnées 0 à 1),
        # donc ce cadre englobe aussi les étiquettes en biais de l'axe des
        # x (ex: "Très satisfait" à 20°) qui débordent de la seule zone de
        # tracé — bbox_inches="tight" recadre ensuite sur ce cadre inclus,
        # jamais en le coupant.
        mpl_fig.patch.set_edgecolor("#B0AAA0")
        mpl_fig.patch.set_linewidth(2.5)
        mpl_fig.tight_layout()
        buf = io.BytesIO()
        # pad_inches augmenté (0.15 -> 0.35) — demande explicite : "trop
        # collé, mets de l'espace" — sans lui, bbox_inches="tight" recadre
        # pile sur le contenu (textes, barres) et coupe la bordure
        # elle-même (fig.patch n'est pas comptée dans le calcul du
        # recadrage "tight"), la faisant disparaître au lieu de rester
        # visible autour de tout, étiquettes en biais comprises.
        mpl_fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", pad_inches=0.35)
        plt.close(mpl_fig)
        return buf.getvalue()
    except Exception:
        return None


def figure_to_png_bytes(fig, width: int = 1200, height: int = 700, scale: float = 2.0,
                         prefer_matplotlib: bool = True, show_title: bool = True) -> bytes | None:
    """Convertit toute figure Plotly en PNG pour les rapports PDF/Word.
    Essaie d'abord un rendu matplotlib pur CPU (sans mini-navigateur) pour
    les figures « simples » (Bar/Scatter/Pie à une seule trace — voir
    `_generic_plotly_to_png_matplotlib`), qui couvrent tous les graphiques
    de la galerie du Générateur de rapport. Ne retombe sur kaleido que pour
    les figures plus complexes que ce convertisseur ne gère pas — le
    rapport doit rester généré (sans image) même si aucun moteur de rendu
    n'est disponible.

    `show_title=False` : n'imprime pas le titre DANS l'image (utile quand
    l'appelant affiche déjà ce titre séparément, en texte PDF/Word — évite
    un doublon visuel du même titre, une fois en image, une fois en texte).

    Mis en cache par (figure, taille) : avec plusieurs agences/graphiques,
    un même rapport peut redemander la même image plusieurs fois (ex : le
    PDF puis le Word du même rapport). Sans cache, chaque clic sur un
    format relance TOUTES les conversions d'image depuis zéro."""
    if fig is None:
        return None
    cache_key = (id(fig), width, height, scale, show_title)
    cached = _PNG_CACHE.get(cache_key)
    if cached is not None:
        return cached

    result = None
    if prefer_matplotlib:
        result = _generic_plotly_to_png_matplotlib(fig, width_px=width, height_px=height, show_title=show_title)

    if result is None:
        try:
            import plotly.io as pio
            result = pio.to_image(fig, format="png", width=width, height=height, scale=scale, engine="kaleido")
        except Exception:
            try:
                result = fig.to_image(format="png", width=width, height=height, scale=scale)
            except Exception:
                result = None
    if result is not None:
        _PNG_CACHE[cache_key] = result
    return result


def chart_png(chart: dict, width: int = 1200, height: int = 700, scale: float = 2.0,
               show_title: bool = True) -> bytes | None:
    """Point d'entrée UNIQUE pour obtenir le PNG d'un graphique de rapport,
    quelle que soit sa nature — figure Plotly (figure_to_png_bytes) OU image
    déjà pré-rendue (ex : nuage de mots, produit directement par la
    bibliothèque `wordcloud`, qui ne passe jamais par Plotly/matplotlib).
    Centralise ce choix pour que les 4 points d'appel (PDF, Word, PPTX,
    pré-rendu) n'aient pas chacun à connaître ce détail."""
    if chart.get("prerendered_png") is not None:
        return chart["prerendered_png"]
    return figure_to_png_bytes(chart.get("fig"), width=width, height=height, scale=scale, show_title=show_title)


# Cache mémoire (durée de vie du process serveur) des images déjà rendues.
# Volontairement simple (dict non borné) : la volumétrie réelle (quelques
# dizaines d'images par session de génération de rapport) reste négligeable
# face au gain de ne jamais re-render la même figure deux fois.
_PNG_CACHE: dict[tuple[int, int, int, float, bool], bytes] = {}


def prerender_agency_page_images(pages: list[dict], progress_callback=None) -> None:
    """Pré-rend en une seule passe, AVANT de construire le PDF/Word, tous les
    PNG nécessaires pour TOUS les graphiques de chaque page (`all_charts` —
    aucune restriction, autant de graphiques que le Tableau de bord en
    produit), avec une progression visible — pour que l'attente sur un
    rapport à plusieurs agences ne ressemble plus jamais à un blocage
    silencieux.

    Rendu en matplotlib en priorité (voir `figure_to_png_bytes`), pas
    kaleido : évite le mini-navigateur intégré, dont le coût par image peut
    devenir prohibitif sur un hébergement aux ressources limitées.

    `progress_callback(done, total, label)` est appelé après chaque image,
    si fourni (typiquement pour alimenter un st.progress côté écran).
    N'a pas besoin d'être appelée pour le PowerPoint (graphiques natifs)."""
    todo = [(p, c) for p in pages for c in (p.get("all_charts") or [])]
    total = len(todo)
    for i, (page, chart) in enumerate(todo):
        chart_png(chart, width=1100, height=650, scale=1.5)
        label = page.get("agence", "") if page else ""
        if progress_callback:
            progress_callback(i + 1, total, label)


def _interpretation(entry: dict) -> str:
    """Interprétation d'une analyse, quel que soit son type. Réutilise le
    moteur de commentaire par règles (viz.comments.generate_comment_rules),
    qui couvre déjà TOUS les types de résultats produits par l'application
    — pas de duplication de logique, pas de type oublié."""
    try:
        from viz.comments import generate_comment_rules
        text = generate_comment_rules(entry)
        if text and "non disponible" not in text:
            return text
    except Exception:  # noqa: BLE001
        pass
    return "Résultat intégré au rapport avec les données disponibles pour cette analyse."


def _clean_md(text: str) -> str:
    """Supprime les marqueurs Markdown visibles (gras/italique) sans jamais
    laisser d'astérisques dans le document final."""
    text = re.sub(r"\*\*\*(.*?)\*\*\*", r"\1", text, flags=re.S)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text, flags=re.S)
    text = re.sub(r"(?<!\*)\*(?!\*)(.*?)\*(?!\*)", r"\1", text, flags=re.S)
    text = text.replace("___", "").replace("__", "")
    return text.replace("***", "")


def _add_runs(paragraph, text: str) -> None:
    """Rendu léger du gras Markdown, sans jamais afficher les astérisques."""
    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**") and len(part) > 4:
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part:
            paragraph.add_run(_clean_md(part))


def _parse_table_lines(lines: list[str]) -> list[list[str]]:
    rows = []
    for line in lines:
        if not line.strip().startswith("|"):
            continue
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if cells and all(re.fullmatch(r":?-{3,}:?", c.replace(" ", "")) for c in cells):
            continue
        rows.append(cells)
    return rows


def _add_table_word(document, rows: list[list[str]]):
    if not rows:
        return None
    from docx.shared import Pt
    cols = max(len(r) for r in rows)
    table = document.add_table(rows=len(rows), cols=cols)
    table.style = "Table Grid"
    for i, row in enumerate(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = _clean_md(row[j] if j < len(row) else "")
            for p in cell.paragraphs:
                for run in p.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(9)
                    if i == 0:
                        run.bold = True
    return table


def _add_markdown_content(document, report_text: str) -> None:
    """Ajoute titres, listes, paragraphes et vrais tableaux Word. Word
    paginant les documents automatiquement, tout le texte est conservé
    quelle que soit sa longueur."""
    lines = (report_text or "").splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if line.strip().startswith("|"):
            block = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                block.append(lines[i])
                i += 1
            _add_table_word(document, _parse_table_lines(block))
            document.add_paragraph()
            continue
        if not line.strip():
            document.add_paragraph()
            i += 1
            continue
        m = re.match(r"^(#{1,3})\s+(.*)", line)
        if m:
            document.add_heading(_clean_md(m.group(2)), level=min(len(m.group(1)) + 1, 4))
            i += 1
            continue
        if line.startswith("- ") or line.startswith("* "):
            p = document.add_paragraph(style="List Bullet")
            _add_runs(p, line[2:])
            i += 1
            continue
        if re.match(r"^\d+\.\s+", line):
            p = document.add_paragraph(style="List Number")
            _add_runs(p, re.sub(r"^\d+\.\s+", "", line))
            i += 1
            continue
        p = document.add_paragraph()
        _add_runs(p, line)
        i += 1


def _variable_label(entry: dict) -> str:
    kind = entry.get("type")
    if kind == "numeric":
        return entry.get("column", "")
    if kind == "categorical":
        return entry.get("column", "")
    if kind in ("correlation",):
        pair = entry.get("strongest_pair", ("", ""))
        return " × ".join(str(x) for x in pair if x)
    if kind == "scatter":
        return f"{entry.get('x','')} × {entry.get('y','')}"
    if kind == "grouped_numeric":
        return f"{entry.get('column','')} par {entry.get('group_column','')}"
    if kind == "timeseries":
        return entry.get("column", "")
    if kind == "rate":
        return entry.get("column", "")
    if kind == "group_test":
        return f"{entry.get('value_col','')} par {entry.get('group_col','')}"
    if kind == "association_test":
        return f"{entry.get('col_a','')} × {entry.get('col_b','')}"
    if kind == "period_comparison":
        return entry.get("label", "")
    if kind == "multi_period_comparison":
        return f"{entry.get('indicator_label','')} ({entry.get('date_col','')})"
    return entry.get("_source_page") or ""


def _summary_rows(analysis_log: list[dict] | None) -> list[list[str]]:
    """Une ligne par analyse de la session, SANS exception — c'est ce
    tableau qui garantit que le rapport couvre bien tous les résultats,
    même ceux qui n'ont pas de graphique associé."""
    rows = [["Analyse", "Variable(s)", "Résultat principal"]]
    for e in analysis_log or []:
        source = e.get("_source_page") or "Session"
        rows.append([source, _variable_label(e), _interpretation(e)])
    return rows


def _split_markdown_sections(report_text: str) -> list[tuple[str, list[str]]]:
    """Découpe le texte Markdown du rapport en sections (titre, lignes),
    pour un rendu PowerPoint fidèle à la structure d'origine plutôt qu'une
    troncature arbitraire."""
    sections: list[tuple[str, list[str]]] = []
    current_title = "Résumé exécutif"
    current_lines: list[str] = []
    for raw in (report_text or "").splitlines():
        line = raw.rstrip()
        m = re.match(r"^#{1,3}\s+(.*)", line)
        if m:
            if current_lines:
                sections.append((current_title, current_lines))
            current_title = _clean_md(m.group(1)) or "Suite"
            current_lines = []
        elif line.strip():
            current_lines.append(_clean_md(line))
    if current_lines:
        sections.append((current_title, current_lines))
    return sections or [("Résumé exécutif", ["Rapport généré à partir des analyses de la session."])]


def _chunk_lines(lines: list[str], max_lines: int = PPTX_MAX_LINES_PER_SLIDE,
                  max_chars: int = PPTX_MAX_CHARS_PER_SLIDE) -> list[list[str]]:
    """Répartit une liste de lignes en groupes qui tiennent sur une seule
    diapositive (ni trop de lignes, ni trop de caractères cumulés) — c'est
    ce mécanisme qui garantit qu'aucune diapositive n'est saturée."""
    chunks: list[list[str]] = []
    current: list[str] = []
    current_chars = 0
    for line in lines:
        if current and (len(current) >= max_lines or current_chars + len(line) > max_chars):
            chunks.append(current)
            current, current_chars = [], 0
        current.append(line)
        current_chars += len(line)
    if current:
        chunks.append(current)
    return chunks or [[]]


def build_word_report(report_text: str, analysis_log: list[dict] | None = None,
                       meta: dict | None = None) -> tuple[bool, bytes | str]:
    """Construit le rapport Word complet : texte rédigé par l'analyste IA,
    tableau de synthèse listant TOUTES les analyses de la session, puis une
    section détaillée avec, pour chaque analyse, son interprétation et son
    graphique quand il existe."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt, RGBColor
    except ImportError:
        return False, "Le package 'python-docx' n'est pas installé (pip install python-docx)."

    try:
        analysis_log = analysis_log or []
        document = Document()
        section = document.sections[0]
        section.top_margin = Inches(.65)
        section.bottom_margin = Inches(.65)
        section.left_margin = Inches(.7)
        section.right_margin = Inches(.7)

        title = document.add_heading("CIE Analytics — Rapport d'analyse", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in title.runs:
            r.font.color.rgb = RGBColor.from_string(ORANGE)
            r.bold = True

        sub = document.add_paragraph(f"Direction Marketing — généré le {datetime.now():%d/%m/%Y à %H:%M}")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in sub.runs:
            r.font.color.rgb = RGBColor.from_string(BLACK)
            r.font.size = Pt(10)

        if meta:
            info = document.add_paragraph()
            info.alignment = WD_ALIGN_PARAGRAPH.CENTER
            bits = []
            if meta.get("lignes_totales") is not None:
                bits.append(f"Lignes analysées : {meta['lignes_totales']:,}".replace(",", " "))
            if meta.get("colonnes"):
                bits.append(f"Colonnes : {len(meta['colonnes'])}")
            if analysis_log:
                bits.append(f"Analyses réalisées : {len(analysis_log)}")
            info.add_run("   |   ".join(bits)).bold = True

        document.add_page_break()
        _add_markdown_content(document, report_text)

        rows = _summary_rows(analysis_log)
        if len(rows) > 1:
            document.add_page_break()
            document.add_heading("Tableau de synthèse — toutes les analyses de la session", 1)
            _add_table_word(document, rows)

        if analysis_log:
            document.add_page_break()
            document.add_heading("Détail des analyses et graphiques", 1)
            n_images = sum(1 for e in analysis_log if e.get("_chart_image"))
            document.add_paragraph(
                f"{len(analysis_log)} analyse(s) réalisée(s) pendant cette session, "
                f"dont {n_images} avec graphique associé."
            )
            for idx, e in enumerate(analysis_log, 1):
                source = e.get("_source_page") or "Analyse"
                p = document.add_paragraph()
                rr = p.add_run(f"{idx}. {source}")
                rr.bold = True
                rr.font.color.rgb = RGBColor.from_string(ORANGE)

                if e.get("_chart_image"):
                    try:
                        document.add_picture(io.BytesIO(e["_chart_image"]), width=Inches(6.45))
                    except Exception as exc:  # noqa: BLE001
                        document.add_paragraph(f"(Graphique non inséré : {exc})")

                ip = document.add_paragraph()
                ir = ip.add_run("Interprétation : ")
                ir.bold = True
                ir.font.color.rgb = RGBColor.from_string(ORANGE)
                ip.add_run(_clean_md(_interpretation(e)))

        # IMPORTANT : ne jamais réassigner `.text` sur un run qui ne
        # contient pas de texte — c'est précisément le run qui porte
        # l'image d'un graphique inséré via `document.add_picture(...)`
        # (python-docx y insère l'image DANS un run). Réassigner `.text`
        # sur ce run, même avec une chaîne vide, remplace tout son contenu
        # XML (donc l'image) par du texte : les graphiques disparaissaient
        # silencieusement du document final. On ne nettoie donc que les
        # runs qui ont réellement du texte à nettoyer.
        for p in document.paragraphs:
            for r in p.runs:
                if r.text:
                    r.text = _clean_md(r.text)

        buffer = io.BytesIO()
        document.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du document Word : {exc}"


def build_pptx_report(report_text: str, analysis_log: list[dict] | None = None,
                       meta: dict | None = None) -> tuple[bool, bytes | str]:
    """Crée un PowerPoint (fond blanc, texte noir, accents orange/jaune CIE)
    dont AUCUNE diapositive n'est surchargée : le texte du rapport et le
    tableau de synthèse sont répartis sur autant de diapositives que
    nécessaire, chaque graphique a sa propre diapositive, et les analyses
    sans graphique sont regroupées par petits lots plutôt qu'omises."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return False, "Le package 'python-pptx' n'est pas installé (pip install python-pptx)."

    try:
        analysis_log = analysis_log or []
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        OR = RGBColor.from_string(ORANGE)
        YE = RGBColor.from_string(YELLOW)
        BK = RGBColor.from_string(BLACK)
        WH = RGBColor.from_string(WHITE)
        BLANK_LAYOUT = prs.slide_layouts[6]

        def add_header(slide, title, subtitle=None):
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.22))
            bar.fill.solid(); bar.fill.fore_color.rgb = OR; bar.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(.65), Inches(.42), Inches(12), Inches(.65))
            p = tb.text_frame.paragraphs[0]
            p.text = _clean_md(title)
            p.font.size = Pt(26 if not subtitle else 24)
            p.font.bold = True
            p.font.color.rgb = BK
            if subtitle:
                sp = tb.text_frame.add_paragraph()
                sp.text = subtitle
                sp.font.size = Pt(14)
                sp.font.color.rgb = OR

        def add_text_slide(title, body_lines, subtitle=None):
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            add_header(slide, title, subtitle)
            tb = slide.shapes.add_textbox(Inches(.75), Inches(1.5), Inches(11.8), Inches(5.3))
            tf = tb.text_frame
            tf.word_wrap = True
            lines = body_lines or ["(pas de contenu pour cette section)"]
            for j, line in enumerate(lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
                p.font.color.rgb = BK
                p.space_after = Pt(10)
            return slide

        # --- Diapositive de titre ------------------------------------------
        title_slide = prs.slides.add_slide(BLANK_LAYOUT)
        bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = WH; bg.line.fill.background()
        band = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), prs.slide_width, Inches(1.6))
        band.fill.solid(); band.fill.fore_color.rgb = YE; band.line.fill.background()
        tb = title_slide.shapes.add_textbox(Inches(.8), Inches(2.1), Inches(11.8), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "CIE Analytics"
        p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = OR
        p2 = tb.text_frame.add_paragraph()
        p2.text = "Rapport d'analyse — Direction Marketing"
        p2.font.size = Pt(28); p2.font.color.rgb = BK

        # --- Texte du rapport, réparti sur autant de diapositives que nécessaire
        sections = _split_markdown_sections(report_text)
        for title, lines in sections:
            chunks = _chunk_lines(lines)
            n_chunks = len(chunks)
            for i, chunk in enumerate(chunks, 1):
                subtitle = f"Partie {i}/{n_chunks}" if n_chunks > 1 else None
                add_text_slide(title, chunk, subtitle=subtitle)

        # --- Tableau de synthèse, réparti sur autant de diapositives que nécessaire
        rows = _summary_rows(analysis_log)
        if len(rows) > 1:
            header, data_rows = rows[0], rows[1:]
            row_chunks = [
                data_rows[i:i + PPTX_MAX_TABLE_ROWS_PER_SLIDE]
                for i in range(0, len(data_rows), PPTX_MAX_TABLE_ROWS_PER_SLIDE)
            ]
            n_chunks = len(row_chunks)
            for i, chunk in enumerate(row_chunks, 1):
                subtitle = f"Partie {i}/{n_chunks}" if n_chunks > 1 else None
                slide = prs.slides.add_slide(BLANK_LAYOUT)
                add_header(slide, "Tableau de synthèse des analyses", subtitle)
                table_rows = [header] + chunk
                table_shape = slide.shapes.add_table(
                    len(table_rows), 3, Inches(.6), Inches(1.5), Inches(12.1), Inches(5.4)
                )
                table = table_shape.table
                for j, w in enumerate([2.3, 3.0, 6.8]):
                    table.columns[j].width = Inches(w)
                for r_idx, row in enumerate(table_rows):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        text = str(val)
                        cell.text = text if len(text) <= 220 else text[:217] + "…"
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = YE if r_idx == 0 else WH
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(9 if r_idx else 11)
                            p.font.bold = (r_idx == 0)
                            p.font.color.rgb = BK

        # --- Une diapositive par graphique -----------------------------------
        with_chart = [e for e in analysis_log if e.get("_chart_image")]
        for idx, e in enumerate(with_chart, 1):
            slide = prs.slides.add_slide(BLANK_LAYOUT)
            add_header(slide, f"Graphique {idx}/{len(with_chart)} — {e.get('_source_page') or 'Analyse'}")
            try:
                slide.shapes.add_picture(io.BytesIO(e["_chart_image"]), Inches(.7), Inches(1.35), width=Inches(7.9), height=Inches(5.6))
            except Exception:  # noqa: BLE001
                pass
            tb = slide.shapes.add_textbox(Inches(8.85), Inches(1.5), Inches(3.8), Inches(5.2))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Interprétation"
            p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = OR
            p2 = tf.add_paragraph()
            p2.text = _interpretation(e)
            p2.font.size = Pt(14); p2.font.color.rgb = BK
            p2.space_before = Pt(10)

        # --- Analyses sans graphique, regroupées par petits lots ------------
        without_chart = [e for e in analysis_log if not e.get("_chart_image")]
        if without_chart:
            bullet_lines = [f"{e.get('_source_page') or 'Analyse'} — {_interpretation(e)}" for e in without_chart]
            chunks = _chunk_lines(bullet_lines, max_lines=6, max_chars=500)
            n_chunks = len(chunks)
            for i, chunk in enumerate(chunks, 1):
                subtitle = f"Partie {i}/{n_chunks}" if n_chunks > 1 else None
                add_text_slide("Autres résultats (sans graphique)", chunk, subtitle=subtitle)

        buffer = io.BytesIO()
        prs.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PowerPoint : {exc}"


def build_pdf_report(report_text: str, analysis_log: list[dict] | None = None,
                      meta: dict | None = None) -> tuple[bool, bytes | str]:
    """Construit le rapport au format PDF : même contenu que le rapport Word
    (texte rédigé, tableau de synthèse, détail des analyses avec graphiques
    et interprétation), mis en page nativement avec ReportLab — aucune
    dépendance externe (pas de LibreOffice requis sur le serveur)."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            Image as RLImage, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
        )
    except ImportError:
        return False, "Le package 'reportlab' n'est pas installé (pip install reportlab)."

    try:
        analysis_log = analysis_log or []
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer, pagesize=A4,
            topMargin=1.6 * cm, bottomMargin=1.6 * cm, leftMargin=1.8 * cm, rightMargin=1.8 * cm,
        )

        styles = getSampleStyleSheet()
        orange = colors.HexColor(f"#{ORANGE}")
        style_title = ParagraphStyle("CieTitle", parent=styles["Title"], textColor=orange, fontSize=22, spaceAfter=4)
        style_sub = ParagraphStyle("CieSub", parent=styles["Normal"], alignment=1, fontSize=9, textColor=colors.black)
        style_h1 = ParagraphStyle("CieH1", parent=styles["Heading1"], textColor=orange, fontSize=15, spaceBefore=14, spaceAfter=8)
        style_h2 = ParagraphStyle("CieH2", parent=styles["Heading2"], textColor=orange, fontSize=12.5, spaceBefore=10, spaceAfter=6)
        style_body = ParagraphStyle("CieBody", parent=styles["Normal"], fontSize=10, leading=14, spaceAfter=6)
        style_bullet = ParagraphStyle("CieBullet", parent=style_body, leftIndent=14, bulletIndent=4)
        style_label = ParagraphStyle("CieLabel", parent=style_body, textColor=orange, fontName="Helvetica-Bold")

        story = []
        story.append(Paragraph("CIE Analytics — Rapport d'analyse", style_title))
        story.append(Paragraph(f"Direction Marketing — généré le {datetime.now():%d/%m/%Y à %H:%M}", style_sub))

        if meta:
            bits = []
            if meta.get("lignes_totales") is not None:
                bits.append(f"Lignes analysées : {meta['lignes_totales']:,}".replace(",", " "))
            if meta.get("colonnes"):
                bits.append(f"Colonnes : {len(meta['colonnes'])}")
            if analysis_log:
                bits.append(f"Analyses réalisées : {len(analysis_log)}")
            if bits:
                story.append(Spacer(1, 4))
                story.append(Paragraph("   |   ".join(bits), ParagraphStyle("CieMeta", parent=style_sub, fontName="Helvetica-Bold")))

        story.append(PageBreak())

        # --- Corps du texte rédigé (Markdown léger : titres, listes, tableaux) ---
        lines = (report_text or "").splitlines()
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()
            if line.strip().startswith("|"):
                block = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    block.append(lines[i])
                    i += 1
                rows = _parse_table_lines(block)
                if rows:
                    story.append(_pdf_table(rows, colors, TableStyle, Table, orange))
                    story.append(Spacer(1, 8))
                continue
            if not line.strip():
                i += 1
                continue
            m = re.match(r"^(#{1,3})\s+(.*)", line)
            if m:
                level = len(m.group(1))
                story.append(Paragraph(_clean_md(m.group(2)), style_h1 if level == 1 else style_h2))
                i += 1
                continue
            if line.startswith("- ") or line.startswith("* "):
                story.append(Paragraph(f"•  {_clean_md(line[2:])}", style_bullet))
                i += 1
                continue
            if re.match(r"^\d+\.\s+", line):
                story.append(Paragraph(_clean_md(re.sub(r'^\d+\.\s+', '', line)), style_bullet))
                i += 1
                continue
            story.append(Paragraph(_clean_md(line), style_body))
            i += 1

        # --- Tableau de synthèse -------------------------------------------
        rows = _summary_rows(analysis_log)
        if len(rows) > 1:
            story.append(PageBreak())
            story.append(Paragraph("Tableau de synthèse — toutes les analyses de la session", style_h1))
            story.append(_pdf_table(rows, colors, TableStyle, Table, orange))

        # --- Détail des analyses, avec graphique et interprétation ----------
        if analysis_log:
            story.append(PageBreak())
            story.append(Paragraph("Détail des analyses et graphiques", style_h1))
            n_images = sum(1 for e in analysis_log if e.get("_chart_image"))
            story.append(Paragraph(
                f"{len(analysis_log)} analyse(s) réalisée(s) pendant cette session, "
                f"dont {n_images} avec graphique associé.", style_body,
            ))
            for idx, e in enumerate(analysis_log, 1):
                source = e.get("_source_page") or "Analyse"
                story.append(Paragraph(f"{idx}. {source}", style_h2))
                if e.get("_chart_image"):
                    try:
                        img_buf = io.BytesIO(e["_chart_image"])
                        story.append(RLImage(img_buf, width=15.5 * cm, height=9 * cm, kind="proportional"))
                        story.append(Spacer(1, 6))
                    except Exception:  # noqa: BLE001
                        story.append(Paragraph("(Graphique non inséré)", style_body))
                story.append(Paragraph("Interprétation :", style_label))
                story.append(Paragraph(_clean_md(_interpretation(e)), style_body))

        doc.build(story)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PDF : {exc}"


def _pdf_table(rows: list[list[str]], colors, TableStyle, Table, header_color):
    """Petit tableau ReportLab uniforme, réutilisé pour la synthèse et les
    tableaux Markdown du corps du rapport."""
    cleaned = [[_clean_md(str(c)) for c in row] for row in rows]
    t = Table(cleaned, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), header_color),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8.5),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#DDDDDD")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FAFAFA")]),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    return t


# ============================================================================
# Rapport "galerie de graphiques" du Générateur de rapport : chaque item est
# {"title": str, "image_bytes": bytes (PNG), "comment": str} — un graphique,
# son commentaire statisticien juste en dessous, pour chacun des formats.
# ============================================================================

def build_gallery_pdf(items: list[dict], meta: dict | None = None) -> tuple[bool, bytes | str]:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import Image as RLImage, PageBreak, Paragraph, SimpleDocTemplate, Spacer
    except ImportError:
        return False, "Le package 'reportlab' n'est pas installé (pip install reportlab)."

    try:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=1.6 * cm, bottomMargin=1.6 * cm,
                                 leftMargin=1.8 * cm, rightMargin=1.8 * cm)
        styles = getSampleStyleSheet()
        orange = colors.HexColor(f"#{ORANGE}")
        style_title = ParagraphStyle("GTitle", parent=styles["Title"], textColor=colors.black, fontName="Helvetica-Bold", fontSize=17, spaceAfter=4)
        style_sub = ParagraphStyle("GSub", parent=styles["Normal"], alignment=1, fontSize=9, textColor=colors.black)
        style_h2 = ParagraphStyle("GH2", parent=styles["Heading2"], textColor=orange, fontSize=13, spaceBefore=6, spaceAfter=6)
        style_label = ParagraphStyle("GLabel", parent=styles["Normal"], fontSize=9.5, textColor=orange, fontName="Helvetica-Bold", spaceBefore=4)
        style_body = ParagraphStyle("GBody", parent=styles["Normal"], fontSize=10, leading=14, spaceAfter=4)

        story = [
            Paragraph("RAPPORT DE SATISFACTION PAR INDICATEURS DE RETOUR DES CLIENTS APRÈS VISITE EN AGENCE DRABO", style_title),
            Paragraph(f"Direction Marketing — généré le {datetime.now():%d/%m/%Y à %H:%M}", style_sub),
        ]
        if meta and meta.get("agences_scope"):
            story.append(Paragraph(f"Agence(s) : {meta['agences_scope']}", style_sub))
        story.append(PageBreak())

        for idx, item in enumerate(items, 1):
            story.append(Paragraph(f"{idx}. {item['title']}", style_h2))
            if item.get("image_bytes"):
                story.append(RLImage(io.BytesIO(item["image_bytes"]), width=16 * cm, height=9 * cm, kind="proportional"))
                story.append(Spacer(1, 4))
            story.append(Paragraph("Commentaire :", style_label))
            story.append(Paragraph(_clean_md(item.get("comment") or "Commentaire non disponible."), style_body))
            story.append(Spacer(1, 14))

        doc.build(story)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PDF : {exc}"


def build_gallery_word(items: list[dict], meta: dict | None = None) -> tuple[bool, bytes | str]:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, RGBColor
    except ImportError:
        return False, "Le package 'python-docx' n'est pas installé (pip install python-docx)."

    try:
        document = Document()
        section = document.sections[0]
        section.top_margin = Inches(.65); section.bottom_margin = Inches(.65)
        section.left_margin = Inches(.7); section.right_margin = Inches(.7)

        title = document.add_heading("RAPPORT DE SATISFACTION PAR INDICATEURS DE RETOUR DES CLIENTS APRÈS VISITE EN AGENCE DRABO", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in title.runs:
            r.font.color.rgb = RGBColor.from_string("000000")
            r.bold = True

        sub = document.add_paragraph(f"Direction Marketing — généré le {datetime.now():%d/%m/%Y à %H:%M}")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if meta and meta.get("agences_scope"):
            scope_p = document.add_paragraph(f"Agence(s) : {meta['agences_scope']}")
            scope_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        document.add_page_break()

        for idx, item in enumerate(items, 1):
            h = document.add_heading(f"{idx}. {item['title']}", level=2)
            for r in h.runs:
                r.font.color.rgb = RGBColor.from_string(ORANGE)
            if item.get("image_bytes"):
                try:
                    document.add_picture(io.BytesIO(item["image_bytes"]), width=Inches(6.45))
                except Exception:  # noqa: BLE001
                    document.add_paragraph("(Graphique non inséré)")
            lbl = document.add_paragraph()
            lr = lbl.add_run("Commentaire : ")
            lr.bold = True
            lr.font.color.rgb = RGBColor.from_string(ORANGE)
            lbl.add_run(_clean_md(item.get("comment") or "Commentaire non disponible."))
            document.add_paragraph()

        buffer = io.BytesIO()
        document.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du document Word : {exc}"


def build_gallery_pptx(items: list[dict], meta: dict | None = None) -> tuple[bool, bytes | str]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor as PptxRGB
        from pptx.util import Inches, Pt
    except ImportError:
        return False, "Le package 'python-pptx' n'est pas installé (pip install python-pptx)."

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        OR = PptxRGB.from_string(ORANGE)
        BK = PptxRGB.from_string(BLACK)

        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(.8), Inches(2.6), Inches(11.7), Inches(1.6))
        tf = tb.text_frame
        tf.text = "RAPPORT DE SATISFACTION PAR INDICATEURS DE RETOUR DES CLIENTS APRÈS VISITE EN AGENCE DRABO"
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BK
        p2 = tf.add_paragraph()
        p2.text = f"Direction Marketing — généré le {datetime.now():%d/%m/%Y à %H:%M}"
        p2.font.size = Pt(16); p2.font.color.rgb = BK
        if meta and meta.get("agences_scope"):
            p3 = tf.add_paragraph()
            p3.text = f"Agence(s) : {meta['agences_scope']}"
            p3.font.size = Pt(14); p3.font.color.rgb = BK

        for idx, item in enumerate(items, 1):
            slide = prs.slides.add_slide(blank)
            title_box = slide.shapes.add_textbox(Inches(.5), Inches(.25), Inches(12.3), Inches(.7))
            title_box.text_frame.text = f"{idx}. {item['title']}"
            title_box.text_frame.paragraphs[0].font.size = Pt(24)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = OR

            if item.get("image_bytes"):
                try:
                    slide.shapes.add_picture(io.BytesIO(item["image_bytes"]), Inches(1.2), Inches(1.05), width=Inches(10.9))
                except Exception:  # noqa: BLE001
                    pass

            cbox = slide.shapes.add_textbox(Inches(.6), Inches(6.15), Inches(12.1), Inches(1.2))
            ctf = cbox.text_frame
            ctf.word_wrap = True
            lbl = ctf.paragraphs[0]
            lbl.text = "Commentaire : " + (item.get("comment") or "Commentaire non disponible.")
            lbl.font.size = Pt(13)
            lbl.font.color.rgb = BK

        buffer = io.BytesIO()
        prs.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PowerPoint : {exc}"


# ============================================================================
# Modèle de rapport officiel CIE — bannière de couverture + une page par
# agence (KPI colorés + anneau + barres), même habillage que le gabarit
# PowerPoint "Indicateurs de retour des clients après visite en agence".
# Images réduites spécifiquement pour l'export (elles restent grandes et
# nettes dans le Tableau de bord — cette réduction ne s'applique qu'ici).
# ============================================================================

CIE_GREEN = "1F7A3F"

def _kpi_pill_pdf(Table, TableStyle, colors, label: str, value: str, bg_hex: str, text_hex: str = "1a1a1a",
                   width_pt: float = 90):
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.platypus import Paragraph
    styles = getSampleStyleSheet()
    # Taille de police proportionnelle à la largeur réellement disponible —
    # avec beaucoup d'indicateurs affichés en même temps, chaque pastille est
    # plus étroite ; sans cette adaptation, les gros chiffres (20pt fixes)
    # débordaient et se retrouvaient coupés au milieu (régression corrigée ici).
    val_size = 20 if width_pt >= 85 else (16 if width_pt >= 65 else 13)
    lbl_size = 10 if width_pt >= 85 else (8.5 if width_pt >= 65 else 7.5)
    val_style = ParagraphStyle("PillVal", parent=styles["Normal"], alignment=1, fontSize=val_size, leading=val_size + 3,
                                fontName=REPORT_FONT_BOLD, textColor=colors.HexColor(f"#{text_hex}"))
    lbl_style = ParagraphStyle("PillLbl", parent=styles["Normal"], alignment=1, fontSize=lbl_size, leading=lbl_size + 3,
                                fontName=REPORT_FONT, textColor=colors.HexColor(f"#{text_hex}"))
    t = Table([[Paragraph(value, val_style)], [Paragraph(label, lbl_style)]], colWidths=[width_pt])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(f"#{bg_hex}")),
        ("TOPPADDING", (0, 0), (-1, 0), 12), ("BOTTOMPADDING", (0, 0), (-1, 0), 4),
        ("TOPPADDING", (0, 1), (-1, 1), 2), ("BOTTOMPADDING", (0, 1), (-1, 1), 12),
        ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
    ]))
    return t


def build_agency_report_pdf(pages: list[dict], meta: dict | None = None, pill_keys: list[str] | None = None) -> tuple[bool, bytes | str]:
    """Rapport officiel : page de couverture (bannière orange/verte CIE) puis
    une page par agence (KPI colorés, anneau de satisfaction, barres des
    points appréciés, commentaire de synthèse)."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            Image as RLImage, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
        )
    except ImportError:
        return False, "Le package 'reportlab' n'est pas installé (pip install reportlab)."

    try:
        buffer = io.BytesIO()
        page_size = A4  # portrait : bien plus de hauteur disponible que le paysage précédent,
                         # qui forçait chaque graphique agrandi sur sa propre page avec un
                         # grand vide en dessous — le format portrait absorbe naturellement
                         # 2 graphiques par page sans gaspillage d'espace.
        doc = SimpleDocTemplate(buffer, pagesize=page_size, topMargin=REPORT_MARGIN_TOP_CM * cm,
                                 bottomMargin=REPORT_MARGIN_BOTTOM_CM * cm,
                                 leftMargin=REPORT_MARGIN_LEFT_CM * cm, rightMargin=REPORT_MARGIN_RIGHT_CM * cm)
        styles = getSampleStyleSheet()
        # Police de base forcée en Times New Roman partout (cohérence totale
        # du rapport) : les styles créés plus bas via parent=styles["Normal"]/
        # ["Heading1"]/["Heading2"]/["Title"] sans fontName explicite
        # hériteraient sinon de la police par défaut de reportlab (Helvetica).
        for _sname in ("Normal", "BodyText", "Heading1", "Heading2", "Heading3", "Title"):
            if _sname in styles.byName:
                _s = styles[_sname]
                _s.fontName = REPORT_FONT_BOLD if ("Head" in _sname or _sname == "Title") else REPORT_FONT
        orange = colors.HexColor(f"#{ORANGE}")
        green = colors.HexColor(f"#{CIE_GREEN}")
        story = []

        # ---------- Logo CIE, uniquement en tout début de rapport ----------
        # Fichier DÉDIÉ (logo_cie_report.png), volontairement DIFFÉRENT de
        # assets/logo_cie.png (utilisé par LOGO_PATH pour la page de
        # connexion et la barre latérale) — chacun son rôle, modifier l'un
        # ne doit jamais changer l'autre par accident.
        _logo_path = Path(__file__).resolve().parent.parent / "assets" / "logo_cie_report.png"
        story.append(Spacer(1, 1.3 * cm))  # descend le logo, trop collé au bord du haut sinon
        if _logo_path.exists():
            story.append(RLImage(str(_logo_path), width=4 * cm, height=4 * cm * (330 / 605), kind="proportional"))
            story.append(Spacer(1, 0.8 * cm))

        # ---------- Page de couverture (bannière orange/verte) ----------
        cover_title_style = ParagraphStyle("CoverTitle", parent=styles["Title"], fontSize=30,
                                            textColor=colors.white, alignment=1, leading=36)
        cover_sub_style = ParagraphStyle("CoverSub", parent=styles["Title"], fontSize=16,
                                          textColor=colors.white, alignment=1, leading=20)
        period_label = meta.get("periode_globale", "") if meta else ""

        band1 = Table([[Paragraph("PÔLE DISTRIBUTION COMMERCIALISATION", cover_sub_style)]],
                       colWidths=[page_size[0] - 5.5 * cm])
        band1.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), green),
                                    ("TOPPADDING", (0, 0), (-1, -1), 14), ("BOTTOMPADDING", (0, 0), (-1, -1), 14)]))
        band2 = Table([[Paragraph("INDICATEURS DE RETOUR DES CLIENTS<br/>APRÈS VISITE EN AGENCE", cover_title_style)]],
                       colWidths=[page_size[0] - 5.5 * cm])
        band2.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), orange),
                                    ("TOPPADDING", (0, 0), (-1, -1), 26), ("BOTTOMPADDING", (0, 0), (-1, -1), 26)]))

        story.append(Spacer(1, 3 * cm))
        story.append(band1)
        story.append(band2)
        if period_label:
            story.append(Spacer(1, 0.6 * cm))
            story.append(Paragraph(period_label, ParagraphStyle("CoverPeriod", parent=styles["Normal"],
                                                                  alignment=1, fontSize=13, textColor=colors.HexColor("#333333"))))
        story.append(PageBreak())

        # ---------- Une page par agence ----------
        style_h1 = ParagraphStyle("AgH1", parent=styles["Heading1"], fontSize=19, textColor=colors.white)
        style_period = ParagraphStyle("AgPeriod", parent=styles["Normal"], fontSize=12.5, textColor=colors.white, alignment=2)
        style_comment_label = ParagraphStyle("AgCLabel", parent=styles["Normal"], fontSize=11,
                                              textColor=orange, fontName=REPORT_FONT_BOLD)
        # Titres de section numérotés (I., II., III.) — structurent chaque
        # page comme un rapport formel (indicateurs / analyse graphique /
        # verbatims), au lieu d'un enchaînement de blocs sans repère.
        style_section = ParagraphStyle("AgSection", parent=styles["Heading2"], fontSize=15,
                                        textColor=colors.HexColor("#1a1a1a"), fontName=REPORT_FONT_BOLD,
                                        spaceBefore=14, spaceAfter=10,
                                        borderWidth=0, borderPadding=0)
        style_comment = ParagraphStyle("AgComment", parent=styles["Normal"], fontName=REPORT_FONT,
                                        fontSize=REPORT_BODY_SIZE, leading=REPORT_BODY_SIZE * REPORT_LINE_SPACING,
                                        alignment=4)  # 4 = justifié, comme Word
        style_chart_title = ParagraphStyle("ChartTitle", parent=styles["Normal"], fontSize=14,
                                            leading=14 * REPORT_LINE_SPACING, textColor=colors.HexColor("#1a1a1a"),
                                            fontName=REPORT_FONT_BOLD, spaceBefore=4, spaceAfter=8)
        style_chart_comment = ParagraphStyle("ChartComment", parent=styles["Normal"], fontSize=REPORT_BODY_SIZE,
                                              leading=REPORT_BODY_SIZE * REPORT_LINE_SPACING,
                                              textColor=colors.HexColor("#333333"),
                                              fontName=REPORT_FONT_ITALIC, spaceAfter=6, alignment=4)  # justifié

        def _section_title(num_str, text):
            """Ligne de titre 'I. Indicateurs clés' avec un filet orange en
            dessous, pour structurer visuellement chaque grande partie de la
            page — inspiré des rapports académiques/professionnels classiques."""
            return [
                Paragraph(f"{num_str}. {text}", style_section),
                Table([[""]], colWidths=[page_size[0] - 5.5 * cm], rowHeights=[1.4],
                      style=TableStyle([("BACKGROUND", (0, 0), (-1, -1), orange)])),
                Spacer(1, 10),
            ]

        def fmt_pct(v):
            return f"{v:.2f} %" if v is not None else "N/A"

        # Numérotation continue des figures (Figure 1, Figure 2...) sur tout
        # le document, toutes agences confondues.
        figure_counter = [0]

        for page in pages:
            k = page["kpis"]
            periode_txt = ""
            if page.get("periode"):
                d0, d1 = page["periode"]
                periode_txt = f"Période : {d0:%d %b} – {d1:%d %b %Y}"

            # Descend le contenu par rapport à la bordure du haut — agrandi
            # (0.6cm -> 1.1cm), toujours trop collé sinon.
            story.append(Spacer(1, 1.1 * cm))

            # Largeur RÉELLEMENT utilisable (marges du document déduites) —
            # bug corrigé ici : l'ancien calcul (page_size[0]*0.6 + *0.32)
            # dépassait la zone de contenu réelle, faisant déborder le
            # bandeau jusque dans la bordure de page des deux côtés.
            usable_w = page_size[0] - 5.5 * cm
            header = Table([[Paragraph(f"RETOUR CLIENTS {page['agence']}", style_h1),
                              Paragraph(periode_txt, style_period)]],
                            colWidths=[usable_w * 0.65, usable_w * 0.35])
            header.setStyle(TableStyle([("BACKGROUND", (0, 0), (0, 0), colors.HexColor("#2a2a2a")),
                                         ("BACKGROUND", (1, 0), (1, 0), orange),
                                         ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                         ("LEFTPADDING", (0, 0), (0, 0), 14), ("TOPPADDING", (0, 0), (-1, -1), 10),
                                         ("BOTTOMPADDING", (0, 0), (-1, -1), 10), ("RIGHTPADDING", (1, 0), (1, 0), 14)]))
            story.append(header)
            story.append(Spacer(1, 18))

            story.extend(_section_title("I", "Indicateurs clés"))

            pill_specs = [
                ("nbre_repondant", "Nbre Répondant", str(k["nbre_repondant"]), "8B4A1E", "ffffff"),
                ("taux_reponse", "Taux de Réponse", fmt_pct(k["taux_reponse"]), "F0F0F0", "1a1a1a"),
                ("taux_satisfaction", "Taux Satisfaction", fmt_pct(k["taux_satisfaction"]), "C8E6C9", "1a1a1a"),
                ("taux_insatisfaction", "Taux d'Insatisfaction", fmt_pct(k["taux_insatisfaction"]), "BDBDBD", "1a1a1a"),
                ("taux_resolution", "Taux de Résolution", fmt_pct(k["taux_resolution"]), "BBDEFB", "1a1a1a"),
            ] + [(key, label, fmt(k.get(key)), bg, fg) for key, label, fmt, bg, fg in EXTRA_KPI_PILLS]
            active_specs = [(lbl, val, bg, fg) for key, lbl, val, bg, fg in pill_specs
                             if key in (pill_keys if pill_keys is not None else [s[0] for s in pill_specs])]
            # Largeur des pastilles calculée UNE SEULE FOIS, sur la base du
            # nombre maximal de pastilles par rangée (5) — jamais recalculée
            # par rangée. Bug corrigé ici : recalculer par rangée rendait les
            # pastilles d'une dernière rangée incomplète (ex: 2 pastilles au
            # lieu de 5) bien plus LARGES que celles des rangées pleines,
            # exactement l'incohérence de taille remontée.
            MAX_PER_ROW = 5
            usable_width = page_size[0] - 5.5 * cm
            pill_w = usable_width / min(MAX_PER_ROW, max(len(active_specs), 1)) - 8
            for row_start in range(0, len(active_specs), MAX_PER_ROW):
                row_specs = active_specs[row_start:row_start + MAX_PER_ROW]
                row_cells = [_kpi_pill_pdf(Table, TableStyle, colors, lbl, val, bg, fg, width_pt=pill_w)
                             for lbl, val, bg, fg in row_specs]
                pills = Table([row_cells])
                pills.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0,0), (-1,-1), 4), ("RIGHTPADDING", (0,0), (-1,-1), 4)]))
                story.append(pills)
                story.append(Spacer(1, 10))
            if active_specs:
                story.append(Spacer(1, 20))

            # TOUS les graphiques de l'agence, deux par ligne, sans aucune
            # restriction — exactement ce que montre le Tableau de bord pour
            # cette agence (voir build_agency_full_charts). Rendu matplotlib
            # (pas kaleido) via figure_to_png_bytes, déjà mis en cache par
            # prerender_agency_page_images.
            # TOUS les graphiques de l'agence, UN SEUL par ligne, en pleine
            # largeur de page — voir la même correction côté export Word :
            # deux graphiques compressés côte à côte rendaient les étiquettes
            # illisibles à la taille d'impression réelle. Chaque graphique a
            # désormais SON PROPRE commentaire juste en dessous (pas
            # uniquement le commentaire de synthèse en fin de page).
            story.extend(_section_title("II", "Analyse graphique détaillée"))
            style_source = ParagraphStyle("SourceNote", parent=styles["Normal"], fontSize=8.5,
                                           fontName=REPORT_FONT_ITALIC, textColor=colors.HexColor("#888888"),
                                           spaceAfter=16)
            for chart in page.get("all_charts") or []:
                img_bytes = chart_png(chart, width=1500, height=850, scale=2.0, show_title=False)
                if img_bytes:
                    if chart.get("title"):
                        figure_counter[0] += 1
                        story.append(Paragraph(f"Figure {figure_counter[0]}. {chart['title']}", style_chart_title))
                    story.append(RLImage(io.BytesIO(img_bytes), width=15 * cm, height=8.6 * cm, kind="proportional"))
                    if chart.get("comment"):
                        story.append(Spacer(1, 8))
                        story.append(Paragraph(f"› {chart['comment']}", style_chart_comment))
                    story.append(Paragraph("Source : données du questionnaire de satisfaction, CIE Analytics.", style_source))
                    story.append(Spacer(1, 10))

            # Commentaire de synthèse de page (bilan global de l'agence) —
            # volontairement MASQUÉ du rapport final : seuls les commentaires
            # SOUS CHAQUE GRAPHIQUE restent affichés (voir ci-dessus).

            verbatims = page.get("verbatims")
            if verbatims and (verbatims.get("positifs") or verbatims.get("negatifs")):
                story.append(Spacer(1, 10))
                story.extend(_section_title("III", "Verbatims clients"))
                for label, items, color_hex in [("Positifs", verbatims.get("positifs") or [], "1B5E20"),
                                                 ("Négatifs", verbatims.get("negatifs") or [], "B71C1C")]:
                    if not items:
                        continue
                    story.append(Paragraph(label, ParagraphStyle(f"Vb{label}", parent=styles["Normal"],
                                                                  fontSize=12, fontName=REPORT_FONT_BOLD,
                                                                  textColor=colors.HexColor(f"#{color_hex}"), spaceBefore=14)))
                    for txt in items:
                        story.append(Paragraph(f"« {txt} »", ParagraphStyle("VbItem", parent=styles["Normal"],
                                                                             fontSize=11.5, leading=18, spaceBefore=8,
                                                                             leftIndent=10)))

            story.append(PageBreak())

        if story and isinstance(story[-1], PageBreak):
            story.pop()

        # Bordure sur CHAQUE page (PDF et couverture comprise) + numéro de
        # page en bas — dessinés directement sur le canevas de chaque page,
        # indépendamment du contenu (story) : s'applique donc uniformément,
        # même si le contenu change de page en page.
        def _decorate_page(canvas, doc_):
            canvas.saveState()
            canvas.setStrokeColor(colors.HexColor(f"#{ORANGE}"))
            # Bordure DOUBLE, marge encore agrandie par rapport au bord de
            # la page (0.7cm -> 0.95cm — reste sous les marges de contenu
            # 1.0/1.2cm pour ne jamais chevaucher le texte) et trait
            # principal à 3pt exactement, comme demandé.
            outer_margin = 0.95 * cm
            inner_margin = outer_margin + 0.16 * cm
            canvas.setLineWidth(3)
            canvas.rect(outer_margin, outer_margin, doc_.pagesize[0] - 2 * outer_margin, doc_.pagesize[1] - 2 * outer_margin)
            canvas.setLineWidth(1)
            canvas.rect(inner_margin, inner_margin, doc_.pagesize[0] - 2 * inner_margin, doc_.pagesize[1] - 2 * inner_margin)
            canvas.setFont(REPORT_FONT, 9)
            canvas.setFillColor(colors.HexColor("#666666"))
            canvas.drawCentredString(doc_.pagesize[0] / 2, outer_margin - 0.4 * cm, f"Page {doc_.page}")
            canvas.restoreState()

        doc.build(story, onFirstPage=_decorate_page, onLaterPages=_decorate_page)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PDF : {exc}"


def build_agency_report_word(pages: list[dict], meta: dict | None = None, pill_keys: list[str] | None = None) -> tuple[bool, bytes | str]:
    """Équivalent Word du modèle officiel CIE : bannière de couverture,
    puis une page par agence (bandeau + pastilles KPI + anneau + barres +
    commentaire), même habillage que la version PDF."""
    try:
        from docx import Document
        from docx.enum.section import WD_ORIENT
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Cm, Inches, Pt, RGBColor
    except ImportError:
        return False, "Le package 'python-docx' n'est pas installé (pip install python-docx)."

    def _shade_cell(cell, hex_color: str):
        shd = cell._tc.get_or_add_tcPr()
        el = shd.makeelement(qn("w:shd"), {qn("w:val"): "clear", qn("w:color"): "auto", qn("w:fill"): hex_color})
        shd.append(el)

    def _set_col_widths(table, widths_in):
        for row in table.rows:
            for cell, w in zip(row.cells, widths_in):
                cell.width = Inches(w)

    try:
        document = Document()
        section = document.sections[0]
        # Portrait (par défaut) : le paysage forçait chaque graphique agrandi
        # sur sa propre page avec un grand vide en dessous — voir la même
        # correction côté PDF pour le détail du raisonnement.
        # Marges au standard d'un rapport statistique professionnel :
        # haut/bas 2,5 cm, gauche 3 cm (reliure), droite 2,5 cm.
        section.top_margin = Cm(2.5); section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(3.0); section.right_margin = Cm(2.5)
        page_w = section.page_width.inches - (3.0 + 2.5) / 2.54  # marges déduites (cm -> pouces)

        # Police par défaut du document : Times New Roman 12pt, interligne
        # 1,5 — appliqué au style "Normal" dont héritent tous les
        # paragraphes qui ne redéfinissent pas explicitement leur police.
        _normal_style = document.styles["Normal"]
        _normal_style.font.name = "Times New Roman"
        _normal_style.font.size = Pt(12)
        _normal_style.paragraph_format.line_spacing = 1.5
        # Police pour les caractères non-latins (accents, etc.) — sans ce
        # réglage East Asian/complex-script, Word peut retomber sur une
        # police différente pour certains caractères accentués.
        _rpr = _normal_style.element.get_or_add_rPr()
        _rFonts = _rpr.find(qn("w:rFonts"))
        if _rFonts is None:
            _rFonts = OxmlElement("w:rFonts")
            _rpr.append(_rFonts)
        _rFonts.set(qn("w:eastAsia"), "Times New Roman")

        # Bordure de page — sur TOUTES les pages du document (une bordure
        # de section s'applique à la section entière, pas page par page ;
        # comme il n'y a qu'une seule section ici, ça couvre tout le
        # document). python-docx n'a pas d'API haut niveau pour ça : on
        # l'injecte directement dans le XML de la section (w:pgBorders).
        # Style "double" natif Word (bordure à deux traits) — épaisseur
        # exactement 3pt (w:sz en huitièmes de point : 3pt = 24) et espace
        # encore agrandi entre bordure et contenu (28 -> 40) pour respirer
        # davantage, comme demandé.
        _sectPr = section._sectPr
        _pgBorders = OxmlElement("w:pgBorders")
        _pgBorders.set(qn("w:offsetFrom"), "page")
        for _side in ("top", "left", "bottom", "right"):
            _border = OxmlElement(f"w:{_side}")
            _border.set(qn("w:val"), "double")
            _border.set(qn("w:sz"), "24")  # 3pt exactement (huitièmes de point)
            _border.set(qn("w:space"), "31")  # maximum autorisé par Word pour offsetFrom="page"
            _border.set(qn("w:color"), ORANGE)
            _pgBorders.append(_border)
        _sectPr.append(_pgBorders)

        # Numéro de page — en pied de page, sur toutes les pages (le champ
        # PAGE se recalcule automatiquement à l'ouverture dans Word).
        _footer_p = section.footer.paragraphs[0]
        _footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _run = _footer_p.add_run()
        _run.font.size = Pt(9)
        _run.font.color.rgb = RGBColor.from_string("666666")
        _fld_begin = OxmlElement("w:fldChar"); _fld_begin.set(qn("w:fldCharType"), "begin")
        _instr = OxmlElement("w:instrText"); _instr.set(qn("xml:space"), "preserve"); _instr.text = "PAGE"
        _fld_end = OxmlElement("w:fldChar"); _fld_end.set(qn("w:fldCharType"), "end")
        _run._r.append(_fld_begin)
        _run._r.append(_instr)
        _run._r.append(_fld_end)

        green = RGBColor.from_string(CIE_GREEN)
        orange = RGBColor.from_string(ORANGE)

        # ---------- Couverture : logo CIE, uniquement en tout début ----------
        # Même fichier dédié que le PDF (assets/logo_cie_report.png),
        # séparé de assets/logo_cie.png (connexion/barre latérale).
        _logo_path = Path(__file__).resolve().parent.parent / "assets" / "logo_cie_report.png"
        document.add_paragraph()  # descend le logo, trop collé au bord du haut sinon
        if _logo_path.exists():
            _logo_p = document.add_paragraph()
            _logo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _logo_p.add_run().add_picture(str(_logo_path), width=Inches(1.6))

        # ---------- Couverture ----------
        for _ in range(4):
            document.add_paragraph()

        band1 = document.add_table(rows=1, cols=1)
        band1.alignment = WD_TABLE_ALIGNMENT.CENTER
        c = band1.rows[0].cells[0]
        _shade_cell(c, CIE_GREEN)
        c.width = Inches(page_w)
        p = c.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run("PÔLE DISTRIBUTION COMMERCIALISATION"); r.bold = True; r.font.size = Pt(16); r.font.color.rgb = RGBColor.from_string("FFFFFF")

        band2 = document.add_table(rows=1, cols=1)
        band2.alignment = WD_TABLE_ALIGNMENT.CENTER
        c = band2.rows[0].cells[0]
        _shade_cell(c, ORANGE)
        c.width = Inches(page_w)
        p = c.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run("INDICATEURS DE RETOUR DES CLIENTS\nAPRÈS VISITE EN AGENCE"); r.bold = True; r.font.size = Pt(26); r.font.color.rgb = RGBColor.from_string("FFFFFF")

        if meta and meta.get("periode_globale"):
            p = document.add_paragraph(meta["periode_globale"])
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.runs[0].font.size = Pt(13)

        p.runs[0].font.color.rgb = RGBColor.from_string("888888")

        document.add_page_break()

        def fmt_pct(v):
            return f"{v:.2f} %" if v is not None else "N/A"

        pill_defs_all = [
            ("nbre_repondant", "Nbre Répondant", lambda k: str(k["nbre_repondant"]), "8B4A1E", "FFFFFF"),
            ("taux_reponse", "Taux de Réponse", lambda k: fmt_pct(k["taux_reponse"]), "F0F0F0", "1A1A1A"),
            ("taux_satisfaction", "Taux Satisfaction", lambda k: fmt_pct(k["taux_satisfaction"]), "C8E6C9", "1A1A1A"),
            ("taux_insatisfaction", "Taux d'Insatisfaction", lambda k: fmt_pct(k["taux_insatisfaction"]), "BDBDBD", "1A1A1A"),
            ("taux_resolution", "Taux de Résolution", lambda k: fmt_pct(k["taux_resolution"]), "BBDEFB", "1A1A1A"),
        ] + [(key, label, (lambda k, key=key, fmt=fmt: fmt(k.get(key))), bg, fg) for key, label, fmt, bg, fg in EXTRA_KPI_PILLS]
        active_keys = pill_keys if pill_keys is not None else [s[0] for s in pill_defs_all]
        pill_defs = [(lbl, getter, bg, fg) for key, lbl, getter, bg, fg in pill_defs_all if key in active_keys]

        # Numérotation continue des figures (Figure 1, Figure 2...) sur tout
        # le document, toutes agences confondues — liste à un élément pour
        # être modifiable depuis l'intérieur de la boucle ci-dessous (pas de
        # "nonlocal" simple en Python pour un entier capturé par closure).
        figure_counter = [0]

        for page_idx, page in enumerate(pages):
            k = page["kpis"]
            periode_txt = ""
            if page.get("periode"):
                d0, d1 = page["periode"]
                periode_txt = f"Période : {d0:%d %b} – {d1:%d %b %Y}"

            # Descend le contenu par rapport à la bordure du haut — agrandi.
            document.add_paragraph().paragraph_format.space_after = Pt(10)
            document.add_paragraph().paragraph_format.space_after = Pt(10)

            header = document.add_table(rows=1, cols=2)
            header.alignment = WD_TABLE_ALIGNMENT.CENTER
            _set_col_widths(header, [page_w * 0.65, page_w * 0.35])
            c0, c1 = header.rows[0].cells
            _shade_cell(c0, "2A2A2A"); _shade_cell(c1, ORANGE)
            p0 = c0.paragraphs[0]; r0 = p0.add_run(f"RETOUR CLIENTS {page['agence']}"); r0.bold = True; r0.font.size = Pt(18); r0.font.color.rgb = RGBColor.from_string("FFFFFF")
            p1 = c1.paragraphs[0]; p1.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            r1 = p1.add_run(periode_txt); r1.font.size = Pt(11); r1.font.color.rgb = RGBColor.from_string("FFFFFF")

            document.add_paragraph().paragraph_format.space_after = Pt(10)

            if pill_defs:
                # Même correctif que le PDF : jamais plus de 5 pastilles par
                # rangée, largeur calculée UNE SEULE FOIS (pas recalculée par
                # rangée) — sinon une dernière rangée incomplète devenait
                # bien plus large que les rangées pleines (bug corrigé).
                MAX_PER_ROW = 5
                cell_w = page_w / min(MAX_PER_ROW, max(len(pill_defs), 1))
                val_pt = 16 if len(pill_defs) <= 5 else 13
                lbl_pt = 8.5 if len(pill_defs) <= 5 else 7.5
                for row_start in range(0, len(pill_defs), MAX_PER_ROW):
                    row_defs = pill_defs[row_start:row_start + MAX_PER_ROW]
                    n_row = len(row_defs)
                    pills = document.add_table(rows=2, cols=n_row)
                    pills.alignment = WD_TABLE_ALIGNMENT.CENTER
                    _set_col_widths(pills, [cell_w] * n_row)
                    for col, (label, getter, bg, fg) in enumerate(row_defs):
                        val_cell = pills.rows[0].cells[col]
                        lbl_cell = pills.rows[1].cells[col]
                        _shade_cell(val_cell, bg); _shade_cell(lbl_cell, bg)
                        vp = val_cell.paragraphs[0]; vp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        vr = vp.add_run(getter(k)); vr.bold = True; vr.font.size = Pt(val_pt); vr.font.color.rgb = RGBColor.from_string(fg)
                        lp = lbl_cell.paragraphs[0]; lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        lr = lp.add_run(label); lr.font.size = Pt(lbl_pt); lr.font.color.rgb = RGBColor.from_string(fg)
                    document.add_paragraph().paragraph_format.space_after = Pt(8)

            document.add_paragraph().paragraph_format.space_after = Pt(10)

            # TOUS les graphiques de l'agence, UN SEUL par ligne, en pleine
            # largeur de page (voir build_agency_full_charts) — un graphique
            # occupant toute la largeur reste lisible une fois imprimé/exporté ;
            # deux graphiques compressés côte à côte rendaient les étiquettes
            # (ex. « Satisfait ») illisibles à la taille d'impression réelle.
            all_charts = page.get("all_charts") or []
            for chart in all_charts:
                document.add_paragraph().paragraph_format.space_after = Pt(10)
                # show_title=False : le titre matplotlib intégré à l'image
                # utilisait une police différente (DejaVu Sans) de tout le
                # reste du document (Times New Roman) — incohérent. Le titre
                # est maintenant un vrai paragraphe Word, même police que le
                # corps du rapport, numéroté comme une figure.
                img_bytes = chart_png(chart, width=1500, height=850, scale=2.0, show_title=False)
                if chart.get("title"):
                    figure_counter[0] += 1
                    tp = document.add_paragraph(); tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    tr = tp.add_run(f"Figure {figure_counter[0]}. {chart['title']}")
                    tr.bold = True; tr.font.size = Pt(13); tr.font.name = "Times New Roman"
                p = document.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if img_bytes:
                    p.add_run().add_picture(io.BytesIO(img_bytes), width=Inches(page_w * 0.92))
                else:
                    p.add_run("Donnée indisponible")
                if chart.get("comment"):
                    document.add_paragraph().paragraph_format.space_after = Pt(2)
                    cp = document.add_paragraph()
                    cp.paragraph_format.space_after = Pt(4)
                    cp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    cr = cp.add_run(f"💬 {chart['comment']}")
                    cr.italic = True; cr.font.size = Pt(12); cr.font.color.rgb = RGBColor.from_string("333333")
                # Source sous chaque figure, comme demandé pour tout rapport
                # statistique professionnel — sobre, petite taille, grisée.
                sp = document.add_paragraph(); sp.paragraph_format.space_after = Pt(16)
                sr = sp.add_run("Source : données du questionnaire de satisfaction, CIE Analytics.")
                sr.italic = True; sr.font.size = Pt(9); sr.font.color.rgb = RGBColor.from_string("888888")

            # Commentaire de synthèse de page volontairement masqué (voir PDF).

            verbatims = page.get("verbatims")
            if verbatims and (verbatims.get("positifs") or verbatims.get("negatifs")):
                document.add_paragraph().paragraph_format.space_after = Pt(14)
                vtitle = document.add_paragraph()
                vr = vtitle.add_run("Verbatims"); vr.bold = True; vr.font.color.rgb = orange; vr.font.size = Pt(14)
                for label, items, color_hex in [("Positifs", verbatims.get("positifs") or [], "1B5E20"),
                                                 ("Négatifs", verbatims.get("negatifs") or [], "B71C1C")]:
                    if not items:
                        continue
                    lp = document.add_paragraph()
                    lp.paragraph_format.space_before = Pt(10)
                    lp.paragraph_format.space_after = Pt(4)
                    lr = lp.add_run(label); lr.bold = True; lr.font.size = Pt(12.5); lr.font.color.rgb = RGBColor.from_string(color_hex)
                    for txt in items:
                        ip = document.add_paragraph(style=None)
                        ip.paragraph_format.left_indent = Inches(0.2)
                        ip.paragraph_format.space_after = Pt(8)
                        ip.paragraph_format.line_spacing = 1.3
                        ip.add_run(f"« {txt} »").font.size = Pt(12)

            if page_idx < len(pages) - 1:
                document.add_page_break()

        buffer = io.BytesIO()
        document.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du document Word : {exc}"


def build_agency_report_pptx(pages: list[dict], meta: dict | None = None,
                              pill_keys: list[str] | None = None) -> tuple[bool, bytes | str]:
    """Équivalent PowerPoint du modèle officiel CIE, AVEC de vrais graphiques
    natifs PowerPoint (anneau et barres) — pas des images : double-clic dans
    PowerPoint pour les modifier, les recolorer, changer le type, exactement
    comme un graphique que l'on aurait construit soi-même dans l'outil. Les
    valeurs viennent des mêmes chiffres exacts que le PDF/Word (donut_data /
    bar_data), jamais recalculées différemment."""
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor as PptxRGB
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Inches, Pt
    except ImportError:
        return False, "Le package 'python-pptx' n'est pas installé (pip install python-pptx)."

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        OR = PptxRGB.from_string(ORANGE)
        GR = PptxRGB.from_string(CIE_GREEN)
        BK = PptxRGB.from_string(BLACK)
        WH = PptxRGB.from_string("FFFFFF")

        # ---------- Diapositive de couverture (bannière verte/orange) ----------
        slide = prs.slides.add_slide(blank)
        band1 = slide.shapes.add_shape(1, Inches(0), Inches(3.0), prs.slide_width, Inches(0.7))  # 1 = RECTANGLE
        band1.fill.solid(); band1.fill.fore_color.rgb = GR; band1.line.fill.background()
        tf1 = band1.text_frame; tf1.text = "PÔLE DISTRIBUTION COMMERCIALISATION"
        tf1.paragraphs[0].font.size = Pt(16); tf1.paragraphs[0].font.bold = True; tf1.paragraphs[0].font.color.rgb = WH
        tf1.word_wrap = True
        from pptx.enum.text import PP_ALIGN
        tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

        band2 = slide.shapes.add_shape(1, Inches(0), Inches(3.7), prs.slide_width, Inches(1.5))
        band2.fill.solid(); band2.fill.fore_color.rgb = OR; band2.line.fill.background()
        tf2 = band2.text_frame; tf2.text = "INDICATEURS DE RETOUR DES CLIENTS\nAPRÈS VISITE EN AGENCE"
        tf2.word_wrap = True
        for p in tf2.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WH; p.alignment = PP_ALIGN.CENTER

        if meta and meta.get("periode_globale"):
            sub = slide.shapes.add_textbox(Inches(0), Inches(5.4), prs.slide_width, Inches(0.5))
            sp = sub.text_frame.paragraphs[0]; sp.text = meta["periode_globale"]
            sp.font.size = Pt(15); sp.font.color.rgb = BK; sp.alignment = PP_ALIGN.CENTER


        # ---------- Une diapositive par agence ----------
        def fmt_pct(v):
            return f"{v:.2f} %" if v is not None else "N/A"

        pill_specs_all = [
            ("nbre_repondant", "Nbre Répondant", lambda k: str(k["nbre_repondant"]), "8B4A1E", "FFFFFF"),
            ("taux_reponse", "Taux de Réponse", lambda k: fmt_pct(k["taux_reponse"]), "F0F0F0", "1A1A1A"),
            ("taux_satisfaction", "Taux Satisfaction", lambda k: fmt_pct(k["taux_satisfaction"]), "C8E6C9", "1A1A1A"),
            ("taux_insatisfaction", "Taux d'Insatisfaction", lambda k: fmt_pct(k["taux_insatisfaction"]), "BDBDBD", "1A1A1A"),
            ("taux_resolution", "Taux de Résolution", lambda k: fmt_pct(k["taux_resolution"]), "BBDEFB", "1A1A1A"),
        ] + [(key, label, (lambda k, key=key, fmt=fmt: fmt(k.get(key))), bg, fg) for key, label, fmt, bg, fg in EXTRA_KPI_PILLS]
        active_keys = pill_keys if pill_keys is not None else [s[0] for s in pill_specs_all]
        pill_specs = [s for s in pill_specs_all if s[0] in active_keys]

        for page in pages:
            k = page["kpis"]
            slide = prs.slides.add_slide(blank)

            periode_txt = ""
            if page.get("periode"):
                d0, d1 = page["periode"]
                periode_txt = f"Période : {d0:%d %b} – {d1:%d %b %Y}"

            # Bandeau titre (noir + orange)
            head_left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(9), Inches(0.85))
            head_left.fill.solid(); head_left.fill.fore_color.rgb = PptxRGB.from_string("2A2A2A"); head_left.line.fill.background()
            hlp = head_left.text_frame.paragraphs[0]; hlp.text = f"RETOUR CLIENTS {page['agence']}"
            hlp.font.size = Pt(22); hlp.font.bold = True; hlp.font.color.rgb = WH
            head_left.text_frame.margin_left = Inches(0.2)

            head_right = slide.shapes.add_shape(1, Inches(9), Inches(0), Inches(4.333), Inches(0.85))
            head_right.fill.solid(); head_right.fill.fore_color.rgb = OR; head_right.line.fill.background()
            hrp = head_right.text_frame.paragraphs[0]; hrp.text = periode_txt
            hrp.font.size = Pt(14); hrp.font.color.rgb = WH; hrp.alignment = PP_ALIGN.RIGHT
            head_right.text_frame.margin_right = Inches(0.2)

            # Pastilles KPI (formes colorées avec texte, pas des images) —
            # même règle que PDF/Word : max 5 par rangée, largeur calculée
            # UNE SEULE FOIS pour que toutes les pastilles (même sur une
            # dernière rangée incomplète) restent strictement identiques.
            MAX_PER_ROW = 5
            n_pills = len(pill_specs)
            if n_pills:
                pill_w = 12.1 / min(MAX_PER_ROW, n_pills)
                for i, (key, label, getter, bg, fg) in enumerate(pill_specs):
                    row, col = divmod(i, MAX_PER_ROW)
                    x = Inches(0.4 + col * pill_w)
                    y = Inches(1.05 + row * 1.05)
                    shp = slide.shapes.add_shape(1, x, y, Inches(pill_w - 0.15), Inches(0.95))
                    shp.fill.solid(); shp.fill.fore_color.rgb = PptxRGB.from_string(bg); shp.line.fill.background()
                    tf = shp.text_frame; tf.word_wrap = True
                    p1 = tf.paragraphs[0]; p1.text = getter(k); p1.font.size = Pt(18); p1.font.bold = True
                    p1.font.color.rgb = PptxRGB.from_string(fg); p1.alignment = PP_ALIGN.CENTER
                    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(9)
                    p2.font.color.rgb = PptxRGB.from_string(fg); p2.alignment = PP_ALIGN.CENTER

            # Commentaire de synthèse de page volontairement masqué (voir PDF).

            # Chaque graphique de l'agence, sur sa propre diapositive, dans
            # l'ordre de build_agency_full_charts — SANS AUCUNE RESTRICTION,
            # exactement ce que montre le Tableau de bord.
            SLIDE_H_IN = 7.5
            SLIDE_CONTENT_LEFT_IN = 0.8
            SLIDE_CONTENT_WIDTH_IN = 11.5
            TOP_IN = 1.0
            for chart in page.get("all_charts") or []:
                extra_slide = prs.slides.add_slide(blank)
                title_box = extra_slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
                tp = title_box.text_frame.paragraphs[0]
                tp.text = f"{page['agence']} — {chart['title']}"
                tp.font.size = Pt(18); tp.font.bold = True; tp.font.color.rgb = OR
                img_bytes = chart_png(chart, width=1100, height=600, scale=1.5)
                # Réserve la place du commentaire AVANT de dimensionner
                # l'image, sinon un commentaire un peu long débordait sous le
                # bas de la diapositive (l'image était toujours calée à 11.5"
                # de large, sans jamais tenir compte de la place restante).
                comment_reserved_in = 1.0 if chart.get("comment") else 0.15
                max_pic_height_in = SLIDE_H_IN - TOP_IN - comment_reserved_in
                pic_bottom_in = TOP_IN
                if img_bytes:
                    from PIL import Image as _PILImage
                    with _PILImage.open(io.BytesIO(img_bytes)) as _im:
                        aspect = _im.width / _im.height  # largeur / hauteur
                    pic_width_in = min(SLIDE_CONTENT_WIDTH_IN, max_pic_height_in * aspect)
                    pic_height_in = pic_width_in / aspect
                    pic_left_in = SLIDE_CONTENT_LEFT_IN + (SLIDE_CONTENT_WIDTH_IN - pic_width_in) / 2
                    extra_slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(pic_left_in), Inches(TOP_IN),
                                                    width=Inches(pic_width_in), height=Inches(pic_height_in))
                    pic_bottom_in = TOP_IN + pic_height_in
                if chart.get("comment"):
                    cbox = extra_slide.shapes.add_textbox(
                        Inches(SLIDE_CONTENT_LEFT_IN), Inches(pic_bottom_in + 0.1),
                        Inches(SLIDE_CONTENT_WIDTH_IN), Inches(SLIDE_H_IN - pic_bottom_in - 0.15),
                    )
                    ctf = cbox.text_frame; ctf.word_wrap = True
                    cp = ctf.paragraphs[0]
                    cp.text = f"💬 {chart['comment']}"
                    cp.font.size = Pt(12); cp.font.italic = True

            # Verbatims : une diapositive dédiée, positifs à gauche / négatifs à droite.
            verbatims = page.get("verbatims")
            if verbatims and (verbatims.get("positifs") or verbatims.get("negatifs")):
                vslide = prs.slides.add_slide(blank)
                vtitle = vslide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
                vtp = vtitle.text_frame.paragraphs[0]
                vtp.text = f"{page['agence']} — Verbatims"
                vtp.font.size = Pt(18); vtp.font.bold = True; vtp.font.color.rgb = OR

                for side_x, label, items, color_hex in [
                    (Inches(0.4), "Positifs", verbatims.get("positifs") or [], "1B5E20"),
                    (Inches(6.8), "Négatifs", verbatims.get("negatifs") or [], "B71C1C"),
                ]:
                    box = vslide.shapes.add_textbox(side_x, Inches(1.0), Inches(6.0), Inches(6.0))
                    tf = box.text_frame; tf.word_wrap = True
                    p0 = tf.paragraphs[0]; p0.text = label
                    p0.font.size = Pt(14); p0.font.bold = True; p0.font.color.rgb = PptxRGB.from_string(color_hex)
                    if not items:
                        p_empty = tf.add_paragraph(); p_empty.text = "Aucun verbatim."
                        p_empty.font.size = Pt(10); p_empty.font.italic = True
                    for txt in items:
                        p = tf.add_paragraph()
                        p.text = f"« {txt} »"
                        p.font.size = Pt(13)
                        p.space_after = Pt(8)

        buffer = io.BytesIO()
        prs.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PowerPoint : {exc}"


# ============================================================================
# Rapport pour données GÉNÉRIQUES (import "Fichier brut quelconque") — même
# moteur de rendu et même habillage visuel (bordure double 3pt, pastilles
# colorées, commentaires) que le rapport CIE, mais SANS notion d'agence : un
# seul jeu d'indicateurs globaux (Lignes, Colonnes, Quanti, Quali...) au lieu
# d'une page par agence, et les graphiques sont ceux RÉELLEMENT configurés
# dans le Tableau de bord générique (variable/opération/type au choix dans
# chaque emplacement) — reconstruits ici à partir de `chartDataSummaries`
# (labels + valeurs déjà calculés côté JS, reçus via le pont Streamlit
# Components), jamais recalculés indépendamment.
# ============================================================================

_GENERIC_KPI_SPECS = [
    ("lignes", "Lignes", "8B4A1E", "ffffff"),
    ("colonnes", "Colonnes", "F0F0F0", "1a1a1a"),
    ("quanti", "Variables quantitatives", "C8E6C9", "1a1a1a"),
    ("quali", "Variables qualitatives", "BBDEFB", "1a1a1a"),
    ("valeurs_manquantes_pct", "Valeurs manquantes (%)", "BDBDBD", "1a1a1a"),
    ("lignes_dupliquees", "Lignes dupliquées", "FFCCBC", "1a1a1a"),
]


def _generic_chart_figure(chart: dict, chart_type: str | None = None):
    """Reconstruit une figure Plotly à partir d'un résumé chiffré déjà
    calculé côté JS (chartDataSummaries) — labels et valeurs identiques à
    ce qui est affiché à l'écran, jamais recalculés.

    `chart_type` (alias de `chart["chart_type"]`, le paramètre explicite
    reste accepté pour compatibilité) : le VRAI type actuellement affiché à
    l'écran pour ce graphique précis — barre-v/barre-h/camembert/anneau/
    courbe/boite/boite_precalc/scatter/groupedbar. Avant, seuls les 5 types
    simples à une seule série étaient gérés, et uniquement via le
    sélecteur de type générique (pas le sélecteur "Type" propre à chaque
    emplacement, le principal) — un nuage de points ou une boîte à
    moustaches disparaissait purement et simplement du rapport (aucun
    label transmis pour ces cas), et un croisement à plusieurs séries
    perdait toutes ses séries sauf la première. Corrigé ici : chaque type
    réellement utilisable dans le Tableau de bord générique a maintenant
    son propre chemin de reconstruction fidèle."""
    import plotly.graph_objects as go
    from viz.charts import apply_readable_style
    from viz.report_charts import apply_chart_type

    vtype = chart_type or chart.get("chart_type")
    labels = chart.get("labels") or []
    datasets = chart.get("datasets") or []
    axis_titles = chart.get("axis_titles") or {}
    title = chart.get("title", "")

    fig = None

    if vtype == "scatter":
        # Nuage de points : un seul dataset, dont chaque valeur est un
        # point {x, y} (pas un simple nombre) — jamais transformable en
        # barres, d'où sa disparition pure et simple avant ce correctif.
        if not datasets:
            return None
        points = datasets[0].get("data") or []
        xs = [p.get("x") for p in points if isinstance(p, dict)]
        ys = [p.get("y") for p in points if isinstance(p, dict)]
        if not xs:
            return None
        fig = go.Figure(go.Scatter(x=xs, y=ys, mode="markers", marker=dict(color=f"#{ORANGE}", size=8)))
        fig.update_layout(xaxis_title=axis_titles.get("x", ""), yaxis_title=axis_titles.get("y", ""))

    elif vtype == "boite":
        # Boîte à moustaches À PARTIR DES VALEURS BRUTES (une boîte par
        # catégorie, ou une seule boîte pour une distribution univariée) —
        # `datasets[0]["data"]` est un tableau de tableaux (les valeurs
        # brutes de chaque boîte), Plotly calcule lui-même les quartiles.
        if not datasets or not labels:
            return None
        raw_groups = datasets[0].get("data") or []
        fig = go.Figure()
        for label, group in zip(labels, raw_groups):
            if isinstance(group, list) and group:
                fig.add_trace(go.Box(y=group, name=str(label), marker_color=f"#{ORANGE}"))
        if not fig.data:
            return None
        fig.update_layout(showlegend=False, xaxis_title=axis_titles.get("x", ""), yaxis_title=axis_titles.get("y", ""))

    elif vtype == "boite_precalc":
        # Boîte à moustaches à partir de statistiques DÉJÀ CALCULÉES
        # (min/Q1/médiane/Q3/max par variable — voir box_stats), utilisée
        # pour le croisement quanti×quanti (2 variables comparées côte à
        # côte). Plotly accepte des quartiles précalculés directement.
        box_stats = chart.get("box_stats") or []
        if not box_stats or not labels:
            return None
        fig = go.Figure(go.Box(
            x=labels,
            q1=[s.get("q1") for s in box_stats], median=[s.get("median") for s in box_stats],
            q3=[s.get("q3") for s in box_stats],
            lowerfence=[s.get("min") for s in box_stats], upperfence=[s.get("max") for s in box_stats],
            marker_color=f"#{ORANGE}",
        ))
        fig.update_layout(showlegend=False, yaxis_title=axis_titles.get("y", ""))

    elif vtype == "groupedbar":
        # Croisement à PLUSIEURS séries (ex : Ville × Produit) — TOUTES les
        # séries, pas seulement la première (bug corrigé : le rapport ne
        # gardait avant que `datasets[0]`, perdant le reste du croisement).
        if not labels or not datasets:
            return None
        fig = go.Figure()
        for i, ds in enumerate(datasets):
            fig.add_trace(go.Bar(name=str(ds.get("label", f"Série {i+1}")), x=[str(v) for v in labels],
                                  y=ds.get("data") or [], marker_color=PALETTE_HEX[i % len(PALETTE_HEX)]))
        fig.update_layout(barmode="group", xaxis_title=axis_titles.get("x", ""), yaxis_title=axis_titles.get("y", ""))

    else:
        # Cas simple, une seule série : barre-v/barre-h/camembert/anneau/courbe.
        if not labels or not datasets:
            return None
        values = datasets[0].get("data") or []
        if len(values) != len(labels):
            return None
        fig = go.Figure(go.Bar(
            x=[str(v) for v in labels], y=values, marker_color=f"#{ORANGE}",
            text=values, textposition="outside",
        ))
        fig.update_layout(title=title)
        apply_readable_style(fig)
        if vtype and vtype != "barre-v":
            fig = apply_chart_type(fig, vtype)
        if axis_titles and vtype in (None, "barre-v", "barre-h"):
            fig.update_layout(xaxis_title=axis_titles.get("x", ""), yaxis_title=axis_titles.get("y", ""))
        return fig

    fig.update_layout(title=title)
    apply_readable_style(fig)
    return fig


def build_generic_report_pdf(charts: list[dict], dataset_kpis: dict, meta: dict | None = None,
                              verbatims: dict | None = None) -> tuple[bool, bytes | str]:
    """Rapport officiel pour données génériques : couverture + UNE page
    d'indicateurs globaux + tous les graphiques configurés (avec leur
    commentaire jusqu'à 4 phrases), même habillage que le rapport CIE."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            Image as RLImage, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
        )

        buffer = io.BytesIO()
        page_size = A4
        # Marges généreuses (>2cm de tous côtés) entre le texte et la
        # bordure décorative — avant, la marge du document (1-1.2cm) était
        # presque égale à la position de la bordure elle-même (0.95cm), le
        # texte semblait quasiment coller au cadre, illisible/à l'étroit.
        doc = SimpleDocTemplate(buffer, pagesize=page_size, topMargin=REPORT_MARGIN_TOP_CM * cm,
                                 bottomMargin=REPORT_MARGIN_BOTTOM_CM * cm,
                                 leftMargin=REPORT_MARGIN_LEFT_CM * cm, rightMargin=REPORT_MARGIN_RIGHT_CM * cm)
        styles = getSampleStyleSheet()
        # Police de base forcée en Times New Roman partout (cohérence totale
        # du rapport) : les styles créés plus bas via parent=styles["Normal"]/
        # ["Heading1"]/["Heading2"]/["Title"] sans fontName explicite
        # hériteraient sinon de la police par défaut de reportlab (Helvetica).
        for _sname in ("Normal", "BodyText", "Heading1", "Heading2", "Heading3", "Title"):
            if _sname in styles.byName:
                _s = styles[_sname]
                _s.fontName = REPORT_FONT_BOLD if ("Head" in _sname or _sname == "Title") else REPORT_FONT
        orange = colors.HexColor(f"#{ORANGE}")
        green = colors.HexColor(f"#{CIE_GREEN}")
        story = []

        _logo_path = Path(__file__).resolve().parent.parent / "assets" / "logo_cie_report.png"
        story.append(Spacer(1, 1.3 * cm))  # descend le logo, trop collé au bord du haut sinon
        if _logo_path.exists():
            story.append(RLImage(str(_logo_path), width=4 * cm, height=4 * cm * (330 / 605), kind="proportional"))
            story.append(Spacer(1, 0.8 * cm))

        cover_title_style = ParagraphStyle("CoverTitleG", parent=styles["Title"], fontSize=28,
                                            textColor=colors.white, alignment=1, leading=34)
        cover_sub_style = ParagraphStyle("CoverSubG", parent=styles["Title"], fontSize=15,
                                          textColor=colors.white, alignment=1, leading=19)
        titre = (meta or {}).get("titre") or "Analyse statistique du jeu de données importé"
        band1 = Table([[Paragraph("DIRECTION MARKETING — ANALYSE DE DONNÉES", cover_sub_style)]],
                       colWidths=[page_size[0] - 5.5 * cm])
        band1.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), green),
                                    ("TOPPADDING", (0, 0), (-1, -1), 14), ("BOTTOMPADDING", (0, 0), (-1, -1), 14)]))
        band2 = Table([[Paragraph(titre.upper(), cover_title_style)]], colWidths=[page_size[0] - 5.5 * cm])
        band2.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), orange),
                                    ("TOPPADDING", (0, 0), (-1, -1), 26), ("BOTTOMPADDING", (0, 0), (-1, -1), 26)]))
        story.append(Spacer(1, 3 * cm))
        story.append(band1)
        story.append(band2)
        story.append(Spacer(1, 0.6 * cm))
        story.append(PageBreak())

        style_h1 = ParagraphStyle("GH1", parent=styles["Heading1"], fontSize=19, textColor=colors.white)
        style_comment_label = ParagraphStyle("GCLabel", parent=styles["Normal"], fontSize=11,
                                              textColor=orange, fontName=REPORT_FONT_BOLD)
        style_chart_title = ParagraphStyle("GChartTitle", parent=styles["Normal"], fontSize=13.5,
                                            leading=17, textColor=colors.HexColor("#1a1a1a"),
                                            fontName=REPORT_FONT_BOLD, spaceBefore=4, spaceAfter=8)
        style_chart_comment = ParagraphStyle("GChartComment", parent=styles["Normal"], fontSize=11.5,
                                              leading=18, textColor=colors.HexColor("#333333"),
                                              fontName=REPORT_FONT_ITALIC, spaceAfter=6, alignment=4)  # justifié
        style_section = ParagraphStyle("GSection", parent=styles["Heading2"], fontSize=15,
                                        textColor=colors.HexColor("#1a1a1a"), fontName=REPORT_FONT_BOLD,
                                        spaceBefore=14, spaceAfter=10)

        def _section_title(num_str, text):
            return [
                Paragraph(f"{num_str}. {text}", style_section),
                Table([[""]], colWidths=[page_size[0] - 5.5 * cm], rowHeights=[1.4],
                      style=TableStyle([("BACKGROUND", (0, 0), (-1, -1), orange)])),
                Spacer(1, 10),
            ]

        story.append(Spacer(1, 1.1 * cm))
        header = Table([[Paragraph("INDICATEURS DU JEU DE DONNÉES", style_h1)]],
                        colWidths=[page_size[0] - 5.5 * cm])
        header.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#2a2a2a")),
                                     ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 14),
                                     ("TOPPADDING", (0, 0), (-1, -1), 10), ("BOTTOMPADDING", (0, 0), (-1, -1), 10)]))
        story.append(header)
        story.append(Spacer(1, 18))

        story.extend(_section_title("I", "Indicateurs du jeu de données"))

        def _fmt_kpi(key, val):
            if val is None:
                return "N/A"
            if key == "valeurs_manquantes_pct":
                return f"{val} %"
            return str(val)

        pill_cells = [_kpi_pill_pdf(Table, TableStyle, colors, label, _fmt_kpi(key, dataset_kpis.get(key)), bg, fg)
                      for key, label, bg, fg in _GENERIC_KPI_SPECS]
        pills = Table([pill_cells])
        pills.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 4),
                                    ("RIGHTPADDING", (0, 0), (-1, -1), 4)]))
        story.append(pills)
        story.append(Spacer(1, 20))

        story.extend(_section_title("II", "Analyses graphiques"))
        style_stat_value = ParagraphStyle("GStatValue", parent=styles["Normal"], fontSize=30, leading=38,
                                           textColor=orange, fontName=REPORT_FONT_BOLD, alignment=1, spaceAfter=10)
        style_source = ParagraphStyle("GSourceNote", parent=styles["Normal"], fontSize=8.5,
                                       fontName=REPORT_FONT_ITALIC, textColor=colors.HexColor("#888888"),
                                       spaceAfter=16)
        figure_counter = [0]  # numérotation continue des figures (Figure 1, Figure 2...)
        for chart in charts:
            if chart.get("stat_value") is not None:
                # Emplacement "valeur unique" (Moyenne, Somme, Médiane...) —
                # aucune image à générer, juste le grand chiffre affiché à
                # l'écran. Disparaissait entièrement du rapport avant ce
                # correctif (aucun label transmis pour ce cas -> ignoré
                # silencieusement par la reconstruction en graphique).
                if chart.get("title"):
                    story.append(Paragraph(chart["title"], style_chart_title))
                story.append(Paragraph(str(chart["stat_value"]), style_stat_value))
                if chart.get("comment"):
                    story.append(Paragraph(f"› {chart['comment']}", style_chart_comment))
                story.append(Spacer(1, 22))
                continue
            fig = _generic_chart_figure(chart, chart_type=chart.get("chart_type"))
            if fig is None:
                continue
            img_bytes = figure_to_png_bytes(fig, width=1500, height=850, scale=2.0, show_title=False)
            if img_bytes:
                if chart.get("title"):
                    figure_counter[0] += 1
                    story.append(Paragraph(f"Figure {figure_counter[0]}. {chart['title']}", style_chart_title))
                story.append(RLImage(io.BytesIO(img_bytes), width=15 * cm, height=8.6 * cm, kind="proportional"))
                if chart.get("comment"):
                    story.append(Spacer(1, 8))
                    story.append(Paragraph(f"› {chart['comment']}", style_chart_comment))
                story.append(Paragraph("Source : données importées via CIE Analytics.", style_source))
                story.append(Spacer(1, 10))

        if verbatims and (verbatims.get("positifs") or verbatims.get("negatifs")):
            story.extend(_section_title("III", "Verbatims"))
            for label, items, color_hex in [("Positifs", verbatims.get("positifs") or [], "1B5E20"),
                                             ("Négatifs", verbatims.get("negatifs") or [], "B71C1C")]:
                if not items:
                    continue
                story.append(Paragraph(label, ParagraphStyle(f"GVb{label}", parent=styles["Normal"],
                                                               fontSize=12, fontName=REPORT_FONT_BOLD,
                                                               textColor=colors.HexColor(f"#{color_hex}"), spaceBefore=12)))
                for txt in items:
                    story.append(Paragraph(f"« {txt} »", ParagraphStyle("GVbItem", parent=styles["Normal"],
                                                                         fontSize=11.5, leading=16, spaceBefore=7,
                                                                         leftIndent=12, alignment=4)))

        def _decorate_page(canvas, doc_):
            canvas.saveState()
            canvas.setStrokeColor(orange)
            outer_margin = 0.95 * cm
            inner_margin = outer_margin + 0.16 * cm
            canvas.setLineWidth(3)
            canvas.rect(outer_margin, outer_margin, doc_.pagesize[0] - 2 * outer_margin, doc_.pagesize[1] - 2 * outer_margin)
            canvas.setLineWidth(1)
            canvas.rect(inner_margin, inner_margin, doc_.pagesize[0] - 2 * inner_margin, doc_.pagesize[1] - 2 * inner_margin)
            canvas.setFont(REPORT_FONT, 9)
            canvas.setFillColor(colors.HexColor("#666666"))
            # À L'INTÉRIEUR de la bordure (avant : en-dessous, à moitié hors
            # cadre — débordait visuellement sous le trait du bas).
            canvas.drawCentredString(doc_.pagesize[0] / 2, outer_margin + 0.35 * cm, f"Page {doc_.page}")
            canvas.restoreState()

        doc.build(story, onFirstPage=_decorate_page, onLaterPages=_decorate_page)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du PDF : {exc}"


def build_generic_report_word(charts: list[dict], dataset_kpis: dict, meta: dict | None = None,
                               verbatims: dict | None = None) -> tuple[bool, bytes | str]:
    """Équivalent Word de `build_generic_report_pdf` — même contenu, même
    bordure double 3pt, même logique de reconstruction des graphiques."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Cm, Inches, Pt, RGBColor

        def _shade_cell(cell, hex_color: str):
            shd = cell._tc.get_or_add_tcPr()
            el = shd.makeelement(qn("w:shd"), {qn("w:val"): "clear", qn("w:color"): "auto", qn("w:fill"): hex_color})
            shd.append(el)

        def _set_col_widths(table, widths_in):
            for row in table.rows:
                for cell, w in zip(row.cells, widths_in):
                    cell.width = Inches(w)

        document = Document()
        section = document.sections[0]
        # Marges au standard d'un rapport statistique professionnel :
        # haut/bas 2,5 cm, gauche 3 cm (reliure), droite 2,5 cm — la bordure
        # de page double reste nettement à l'intérieur de cette marge.
        section.top_margin = Cm(2.5); section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(3.0); section.right_margin = Cm(2.5)
        page_w = section.page_width.inches - (3.0 + 2.5) / 2.54

        # Police par défaut du document : Times New Roman 12pt, interligne 1,5.
        _normal_style = document.styles["Normal"]
        _normal_style.font.name = "Times New Roman"
        _normal_style.font.size = Pt(12)
        _normal_style.paragraph_format.line_spacing = 1.5
        _rpr = _normal_style.element.get_or_add_rPr()
        _rFonts = _rpr.find(qn("w:rFonts"))
        if _rFonts is None:
            _rFonts = OxmlElement("w:rFonts")
            _rpr.append(_rFonts)
        _rFonts.set(qn("w:eastAsia"), "Times New Roman")

        _sectPr = section._sectPr
        _pgBorders = OxmlElement("w:pgBorders")
        _pgBorders.set(qn("w:offsetFrom"), "page")
        for _side in ("top", "left", "bottom", "right"):
            _border = OxmlElement(f"w:{_side}")
            _border.set(qn("w:val"), "double")
            _border.set(qn("w:sz"), "24")
            _border.set(qn("w:space"), "31")
            _border.set(qn("w:color"), ORANGE)
            _pgBorders.append(_border)
        _sectPr.append(_pgBorders)

        _footer_p = section.footer.paragraphs[0]
        _footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _run = _footer_p.add_run()
        _run.font.size = Pt(9)
        _run.font.color.rgb = RGBColor.from_string("666666")
        _fld_begin = OxmlElement("w:fldChar"); _fld_begin.set(qn("w:fldCharType"), "begin")
        _instr = OxmlElement("w:instrText"); _instr.set(qn("xml:space"), "preserve"); _instr.text = "PAGE"
        _fld_end = OxmlElement("w:fldChar"); _fld_end.set(qn("w:fldCharType"), "end")
        _run._r.append(_fld_begin); _run._r.append(_instr); _run._r.append(_fld_end)

        green = RGBColor.from_string(CIE_GREEN)
        orange = RGBColor.from_string(ORANGE)

        _logo_path = Path(__file__).resolve().parent.parent / "assets" / "logo_cie_report.png"
        document.add_paragraph()  # descend le logo, trop collé au bord du haut sinon
        if _logo_path.exists():
            _logo_p = document.add_paragraph()
            _logo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _logo_p.add_run().add_picture(str(_logo_path), width=Inches(1.6))

        for _ in range(4):
            document.add_paragraph()

        titre = (meta or {}).get("titre") or "Analyse statistique du jeu de données importé"
        band1 = document.add_table(rows=1, cols=1)
        _set_col_widths(band1, [page_w])
        c0 = band1.rows[0].cells[0]; _shade_cell(c0, CIE_GREEN)
        p0 = c0.paragraphs[0]; p0.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r0 = p0.add_run("DIRECTION MARKETING — ANALYSE DE DONNÉES"); r0.bold = True; r0.font.size = Pt(13); r0.font.color.rgb = RGBColor.from_string("FFFFFF")

        band2 = document.add_table(rows=1, cols=1)
        _set_col_widths(band2, [page_w])
        c1 = band2.rows[0].cells[0]; _shade_cell(c1, ORANGE)
        p1 = c1.paragraphs[0]; p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r1 = p1.add_run(titre.upper()); r1.bold = True; r1.font.size = Pt(22); r1.font.color.rgb = RGBColor.from_string("FFFFFF")

        document.add_paragraph()
        document.add_page_break()

        header = document.add_table(rows=1, cols=1)
        _set_col_widths(header, [page_w])
        hc = header.rows[0].cells[0]; _shade_cell(hc, "2A2A2A")
        hp = hc.paragraphs[0]; hr = hp.add_run("INDICATEURS DU JEU DE DONNÉES"); hr.bold = True; hr.font.size = Pt(18); hr.font.color.rgb = RGBColor.from_string("FFFFFF")

        document.add_paragraph().paragraph_format.space_after = Pt(10)

        def _fmt_kpi(key, val):
            if val is None:
                return "N/A"
            if key == "valeurs_manquantes_pct":
                return f"{val} %"
            return str(val)

        n_pills = len(_GENERIC_KPI_SPECS)
        cell_w = page_w / n_pills
        pills = document.add_table(rows=2, cols=n_pills)
        pills.alignment = 1
        _set_col_widths(pills, [cell_w] * n_pills)
        for col, (key, label, bg, fg) in enumerate(_GENERIC_KPI_SPECS):
            val_cell = pills.rows[0].cells[col]; lbl_cell = pills.rows[1].cells[col]
            _shade_cell(val_cell, bg); _shade_cell(lbl_cell, bg)
            vp = val_cell.paragraphs[0]; vp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            vr = vp.add_run(_fmt_kpi(key, dataset_kpis.get(key))); vr.bold = True; vr.font.size = Pt(15); vr.font.color.rgb = RGBColor.from_string(fg)
            lp = lbl_cell.paragraphs[0]; lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            lr = lp.add_run(label); lr.font.size = Pt(8); lr.font.color.rgb = RGBColor.from_string(fg)

        document.add_paragraph().paragraph_format.space_after = Pt(10)

        figure_counter = [0]  # numérotation continue des figures
        for chart in charts:
            if chart.get("stat_value") is not None:
                # Emplacement "valeur unique" (Moyenne, Somme, Médiane...) —
                # disparaissait entièrement du rapport avant ce correctif.
                if chart.get("title"):
                    tp = document.add_paragraph(); tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    tr = tp.add_run(chart["title"]); tr.bold = True; tr.font.size = Pt(13); tr.font.name = "Times New Roman"
                vp = document.add_paragraph(); vp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                vr = vp.add_run(str(chart["stat_value"])); vr.bold = True; vr.font.size = Pt(30)
                vr.font.color.rgb = orange
                if chart.get("comment"):
                    cp = document.add_paragraph(); cp.paragraph_format.space_after = Pt(16)
                    cp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    cr = cp.add_run(f"💬 {chart['comment']}")
                    cr.italic = True; cr.font.size = Pt(12); cr.font.color.rgb = RGBColor.from_string("333333")
                continue
            fig = _generic_chart_figure(chart, chart_type=chart.get("chart_type"))
            if fig is None:
                continue
            document.add_paragraph().paragraph_format.space_after = Pt(10)
            if chart.get("title"):
                figure_counter[0] += 1
                tp = document.add_paragraph(); tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                tr = tp.add_run(f"Figure {figure_counter[0]}. {chart['title']}")
                tr.bold = True; tr.font.size = Pt(13); tr.font.name = "Times New Roman"
            img_bytes = figure_to_png_bytes(fig, width=1500, height=850, scale=2.0, show_title=False)
            p = document.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if img_bytes:
                p.add_run().add_picture(io.BytesIO(img_bytes), width=Inches(page_w * 0.92))
            if chart.get("comment"):
                document.add_paragraph().paragraph_format.space_after = Pt(2)
                cp = document.add_paragraph(); cp.paragraph_format.space_after = Pt(4)
                cp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                cr = cp.add_run(f"💬 {chart['comment']}")
                cr.italic = True; cr.font.size = Pt(12); cr.font.color.rgb = RGBColor.from_string("333333")
            sp = document.add_paragraph(); sp.paragraph_format.space_after = Pt(16)
            sr = sp.add_run("Source : données importées via CIE Analytics.")
            sr.italic = True; sr.font.size = Pt(9); sr.font.color.rgb = RGBColor.from_string("888888")

        if verbatims and (verbatims.get("positifs") or verbatims.get("negatifs")):
            document.add_page_break()
            vheader = document.add_table(rows=1, cols=1)
            _set_col_widths(vheader, [page_w])
            vhc = vheader.rows[0].cells[0]; _shade_cell(vhc, "2A2A2A")
            vhp = vhc.paragraphs[0]; vhr = vhp.add_run("VERBATIMS"); vhr.bold = True; vhr.font.size = Pt(18); vhr.font.color.rgb = RGBColor.from_string("FFFFFF")
            document.add_paragraph().paragraph_format.space_after = Pt(10)
            for label, items, color_hex in [("Positifs", verbatims.get("positifs") or [], "1B5E20"),
                                             ("Négatifs", verbatims.get("negatifs") or [], "B71C1C")]:
                if not items:
                    continue
                lp = document.add_paragraph()
                lr = lp.add_run(label); lr.bold = True; lr.font.size = Pt(13); lr.font.color.rgb = RGBColor.from_string(color_hex)
                for txt in items:
                    ip = document.add_paragraph(); ip.paragraph_format.left_indent = Inches(.25)
                    ip.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    ir = ip.add_run(f"« {txt} »"); ir.font.size = Pt(11.5)

        buffer = io.BytesIO()
        document.save(buffer)
        return True, buffer.getvalue()
    except Exception as exc:  # noqa: BLE001
        return False, f"Erreur lors de la génération du document Word : {exc}"
